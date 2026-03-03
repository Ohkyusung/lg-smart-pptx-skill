#!/usr/bin/env python3
"""
LG Group PPTX Builder (LG Smart Font Edition)
LG 그룹 브랜드 가이드라인을 따르는 프레젠테이션 생성 라이브러리
LG스마트체(LG Smart) 폰트를 기본 적용합니다.

Usage:
    from lg_pptx_builder import LGPresentation

    prs = LGPresentation()
    prs.add_cover("프로젝트 제목", subtitle="팀명", date="2025.10.23")
    prs.add_toc([("Summary", []), ("시스템 소개", ["항목1", "항목2"])])
    prs.add_content("1.1 시스템 개요", section="Summary", bullets=["내용1", "내용2"])
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─────────────────────────────────────────────
# Design Tokens
# ─────────────────────────────────────────────

class LGColors:
    """LG Group brand color palette"""
    RED = RGBColor(0xA5, 0x00, 0x34)           # Primary - LG RED
    BLACK = RGBColor(0x00, 0x00, 0x00)          # Text primary
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)      # Text secondary
    MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)    # Text tertiary / section labels
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)     # Surface / content box bg
    BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)     # Borders / dividers
    CHARCOAL = RGBColor(0x3C, 0x3C, 0x3C)       # Header bars / timeline
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # Background
    GREEN = RGBColor(0x2E, 0x7D, 0x32)          # Future / planned items
    ORANGE = RGBColor(0xD4, 0x76, 0x0A)         # Highlight items


class LGTypography:
    """LG typography scale"""
    COVER_TITLE = Pt(32)
    COVER_SUBTITLE = Pt(14)
    SECTION_TITLE = Pt(28)       # Slide title (Bold) - tighter
    SUBTITLE = Pt(16)            # Slide subtitle (SemiBold) - tighter
    BODY_TITLE = Pt(16)
    BODY = Pt(12)                # All detail/body text (Regular) - denser
    BODY_SMALL = Pt(11)
    TABLE_HEADER = Pt(10)
    TABLE_BODY = Pt(9)
    CAPTION = Pt(9)
    TOC_TITLE = Pt(28)
    TOC_ITEM = Pt(16)
    TOC_SUBITEM = Pt(13)


class LGDimensions:
    """Slide layout dimensions (16:9, 13.333 x 7.5 inches)"""
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins (report-dense layout)
    LEFT_MARGIN = Cm(0.8)
    RIGHT_MARGIN = Cm(0.5)
    TOP_MARGIN = Cm(0.5)
    BOTTOM_MARGIN = Cm(0.5)

    # Accent bar
    ACCENT_BAR_WIDTH = Cm(0.3)
    ACCENT_BAR_LEFT = Cm(0)

    # Content area (after accent bar)
    CONTENT_LEFT = Cm(1.5)
    CONTENT_TOP = Cm(0.5)
    CONTENT_WIDTH = Cm(31.5)  # ~full width minus margins

    # L-bracket
    BRACKET_ARM_LENGTH = Cm(2.5)
    BRACKET_THICKNESS = Cm(0.4)


def _make_oxml_element(tag, **attribs):
    """Create an OxmlElement with attributes."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    element = etree.SubElement(etree.Element('dummy', nsmap=nsmap), qn(tag))
    for k, v in attribs.items():
        element.set(k, str(v))
    # Detach from dummy parent
    dummy = element.getparent()
    dummy.remove(element)
    return element


# ─────────────────────────────────────────────
# Main Builder Class
# ─────────────────────────────────────────────

class LGPresentation:
    """
    LG 그룹 브랜드 가이드라인을 따르는 PPTX 프레젠테이션 빌더 (LG Smart 폰트)

    Args:
        font_name: 폰트 패밀리 이름 (기본: "LG Smart")
        font_name_latin: 라틴 폰트 이름 (기본: "LG Smart")
        fallback_font: 폴백 폰트 (기본: "맑은 고딕")
        logo_path: LG 로고 이미지 파일 경로 (선택)
    """

    def __init__(self, font_name="LG Smart", font_name_latin="LG Smart",
                 fallback_font="맑은 고딕", logo_path=None):
        self.font_name = font_name
        self.font_name_latin = font_name_latin
        self.fallback_font = fallback_font
        self.logo_path = logo_path

        # LG Smart uses a standard font family structure:
        #   "LG Smart" (Regular, Bold) — shared family name
        #   "LG Smart Light" (Regular=Light, Bold=SemiBold) — light sub-family
        # PowerPoint distinguishes weights via the bold flag, not family name.
        self.font_regular = font_name            # "LG Smart"
        self.font_semibold = f"{font_name} Light" # "LG Smart Light" (Bold = SemiBold)
        self.font_bold = font_name                # "LG Smart" (with bold flag)

        self.prs = Presentation()
        self.prs.slide_width = LGDimensions.SLIDE_WIDTH
        self.prs.slide_height = LGDimensions.SLIDE_HEIGHT

        # Patch the presentation theme so the default font resolves correctly.
        self._patch_theme_fonts()

    # ─────────────────────────────────────────
    # Font Helpers
    # ─────────────────────────────────────────

    def _patch_theme_fonts(self):
        """
        Rewrite the presentation theme XML so that the default fonts use
        LG Smart.

        LG Smart uses a standard family structure where Regular and Bold
        share the same "LG Smart" family name.  PowerPoint resolves the
        weight via the bold flag, so both majorFont and minorFont point
        to the same family.

        We set:
          - majorFont (headings) → "LG Smart"
          - minorFont (body text) → "LG Smart"
        """
        slide_master = self.prs.slide_masters[0]
        theme_part = None
        for rel in slide_master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                break

        if theme_part is None:
            return

        theme_xml = etree.fromstring(theme_part.blob)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        # Both major and minor use the same family; PowerPoint picks
        # Regular vs Bold via the bold attribute on each run.
        font_map = {
            'a:majorFont': self.font_bold,     # "LG Smart"
            'a:minorFont': self.font_regular,   # "LG Smart"
        }

        for font_group_tag, target_name in font_map.items():
            font_group = theme_xml.find(f'.//a:fontScheme/{font_group_tag}', ns)
            if font_group is None:
                continue

            for child_tag in ['a:latin', 'a:ea', 'a:cs']:
                child = font_group.find(child_tag, ns)
                if child is not None:
                    child.set('typeface', target_name)

            for font_el in font_group.findall('a:font', ns):
                if font_el.get('script') == 'Hang':
                    font_el.set('typeface', target_name)

        theme_part._blob = etree.tostring(
            theme_xml, xml_declaration=True, encoding='UTF-8', standalone=True
        )

    def _set_font(self, run, size=None, bold=False, italic=False, color=None,
                  weight=None):
        """
        Set font properties including East Asian font for Korean text.

        Args:
            weight: Explicit weight override — "bold", "semibold", or "regular".
                    If omitted, falls back to ``bold`` flag (True→bold, False→regular).
        """
        # LG Smart uses a standard family structure:
        #   "LG Smart" for Regular (bold=False) and Bold (bold=True)
        #   "LG Smart Light" for SemiBold (registered as Bold of Light sub-family)
        if weight == "semibold":
            target_font = self.font_semibold   # "LG Smart Light"
            use_bold = True                     # SemiBold = Bold of Light sub-family
        elif weight == "bold" or (weight is None and bold):
            target_font = self.font_bold       # "LG Smart"
            use_bold = True
        else:
            target_font = self.font_regular    # "LG Smart"
            use_bold = False

        font = run.font
        if size:
            font.size = size
        font.bold = use_bold
        font.italic = italic
        if color:
            font.color.rgb = color

        # Directly manipulate XML for reliable font binding
        rPr = run._r.get_or_add_rPr()

        # Set language attributes on the run properties element
        rPr.set('lang', 'ko-KR')
        rPr.set('altLang', 'en-US')

        # Remove existing font elements to avoid duplicates
        for tag in ['a:latin', 'a:ea', 'a:cs']:
            for existing in rPr.findall(qn(tag)):
                rPr.remove(existing)

        # <a:latin> — Latin font
        latin = _make_oxml_element('a:latin')
        latin.set('typeface', target_font)
        rPr.append(latin)

        # <a:ea> — East Asian font (charset="-127" for CJK binding)
        ea = _make_oxml_element('a:ea')
        ea.set('typeface', target_font)
        ea.set('charset', '-127')
        rPr.append(ea)

        # <a:cs> — Complex Script font
        cs = _make_oxml_element('a:cs')
        cs.set('typeface', target_font)
        cs.set('charset', '-127')
        rPr.append(cs)

    def _add_text(self, text_frame, text, size=None, bold=False, color=None,
                  alignment=None, space_before=None, space_after=None,
                  weight=None):
        """Add a paragraph with formatted text to a text frame."""
        if text_frame.paragraphs and text_frame.paragraphs[0].text == '':
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if alignment:
            p.alignment = alignment
        if space_before:
            p.space_before = space_before
        if space_after:
            p.space_after = space_after

        run = p.add_run()
        run.text = text
        self._set_font(run, size=size, bold=bold, color=color, weight=weight)
        return p

    def _add_textbox(self, slide, left, top, width, height, text="",
                     size=None, bold=False, color=None, alignment=None,
                     word_wrap=True, vertical=None, weight=None):
        """Add a text box with formatted text."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        if vertical:
            tf.vertical_anchor = vertical

        if text:
            self._add_text(tf, text, size=size, bold=bold, color=color,
                          alignment=alignment, weight=weight)

        # Make transparent
        txBox.fill.background()
        txBox.line.fill.background()
        return txBox

    # ─────────────────────────────────────────
    # Visual Element Helpers
    # ─────────────────────────────────────────

    def _add_l_bracket(self, slide, corner="top-left", arm_length=None,
                       thickness=None, color=None):
        """
        Draw an L-bracket corner decoration.

        Args:
            corner: "top-left" or "bottom-right"
            arm_length: Length of bracket arms (default: 2.5cm)
            thickness: Thickness of bracket arms (default: 0.4cm)
            color: Fill color (default: LG RED)
        """
        arm = arm_length or LGDimensions.BRACKET_ARM_LENGTH
        thick = thickness or LGDimensions.BRACKET_THICKNESS
        clr = color or LGColors.RED

        shapes = slide.shapes
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)
        a = int(arm)
        t = int(thick)

        if corner == "top-left":
            # Horizontal arm (top)
            h_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(0.8), Cm(0.8), a, t)
            h_bar.fill.solid()
            h_bar.fill.fore_color.rgb = clr
            h_bar.line.fill.background()

            # Vertical arm (left)
            v_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(0.8), Cm(0.8), t, a)
            v_bar.fill.solid()
            v_bar.fill.fore_color.rgb = clr
            v_bar.line.fill.background()

        elif corner == "bottom-right":
            # Position from bottom-right
            br_x = sw - int(Cm(0.8)) - a
            br_y = sh - int(Cm(0.8)) - a

            # Horizontal arm (bottom)
            h_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     br_x, br_y + a - t, a, t)
            h_bar.fill.solid()
            h_bar.fill.fore_color.rgb = clr
            h_bar.line.fill.background()

            # Vertical arm (right)
            v_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     br_x + a - t, br_y, t, a)
            v_bar.fill.solid()
            v_bar.fill.fore_color.rgb = clr
            v_bar.line.fill.background()

    def _add_accent_bar(self, slide, color=None, top=None):
        """
        Add a small accent block in the top-left, next to the slide title.
        This is a short red rectangle that acts as a title marker,
        matching the LG content slide pattern (see reference images).

        Args:
            color: Fill color (default: LG RED)
            top: Top position (default: 0.7cm, aligned with title)
        """
        clr = color or LGColors.RED
        t = top if top is not None else Cm(0.4)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(0.8), t, Cm(0.4), Cm(1.0)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = clr
        bar.line.fill.background()
        return bar

    def _add_section_indicator(self, slide, section_name, has_dot=True):
        """Add section name label in top-right with optional red dot."""
        # Section name text
        sw = int(LGDimensions.SLIDE_WIDTH)
        txBox = self._add_textbox(
            slide,
            left=sw - Cm(10), top=Cm(0.3),
            width=Cm(9), height=Cm(0.8),
            text=section_name,
            size=LGTypography.BODY_SMALL,
            color=LGColors.MEDIUM_GRAY,
            alignment=PP_ALIGN.RIGHT
        )

        if has_dot:
            # Red dot indicator
            dot_size = Cm(0.35)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                sw - Cm(0.8), Cm(0.5),
                dot_size, dot_size
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = LGColors.RED
            dot.line.fill.background()

        return txBox

    def _add_slide_title(self, slide, title, left=None, top=None):
        """Add slide title text (for content slides, positioned right of accent block)."""
        l = left or Cm(1.5)
        t = top or Cm(0.3)
        return self._add_textbox(
            slide, l, t,
            width=Cm(25), height=Cm(1.2),
            text=title,
            size=Pt(24),
            bold=True,
            color=LGColors.BLACK,
            vertical=MSO_ANCHOR.MIDDLE
        )

    def _add_horizontal_line(self, slide, left, top, width, color=None, height=None):
        """Add a thin horizontal line."""
        clr = color or LGColors.BORDER_GRAY
        h = height or Pt(1)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, h
        )
        line.fill.solid()
        line.fill.fore_color.rgb = clr
        line.line.fill.background()
        return line

    def _add_shadow(self, shape, blur_radius=50800, dist=38100,
                    direction=2700000, color="000000", alpha=40):
        """
        Add a subtle outer drop shadow to a shape for depth/dimension.

        This injects an <a:effectLst> containing an <a:outerShdw> element
        into the shape's spPr (shape properties) XML.

        Args:
            shape: The shape object to add a shadow to
            blur_radius: Shadow blur in EMU (default ~4pt)
            dist: Shadow distance in EMU (default ~3pt)
            direction: Shadow angle in 60000ths of a degree (default 2700000 = 270° = below-right)
            color: Shadow color hex string without '#' (default "000000")
            alpha: Shadow opacity percentage 0-100 (default 40)
        """
        spPr = shape._element.spPr

        # Remove any existing effectLst
        for existing in spPr.findall(qn('a:effectLst')):
            spPr.remove(existing)

        effectLst = _make_oxml_element('a:effectLst')

        outerShdw = _make_oxml_element('a:outerShdw')
        outerShdw.set('blurRad', str(blur_radius))
        outerShdw.set('dist', str(dist))
        outerShdw.set('dir', str(direction))
        outerShdw.set('rotWithShape', '0')

        srgbClr = _make_oxml_element('a:srgbClr')
        srgbClr.set('val', color)

        alphaElem = _make_oxml_element('a:alpha')
        alphaElem.set('val', str(alpha * 1000))  # OOXML uses 1000ths of a percent
        srgbClr.append(alphaElem)

        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)

    def _get_blank_slide(self):
        """Get a blank slide layout and add a new slide."""
        # Try to find blank layout (usually index 6, but may vary)
        layout = self.prs.slide_layouts[6]  # Blank
        return self.prs.slides.add_slide(layout)

    # ─────────────────────────────────────────
    # Table Helpers
    # ─────────────────────────────────────────

    def _set_cell_border(self, cell, color="CCCCCC", width='6350'):
        """Set all four borders on a table cell."""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for edge in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            for existing in tcPr.findall(qn(edge)):
                tcPr.remove(existing)
            ln = _make_oxml_element(edge)
            ln.set('w', str(width))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')

            solidFill = _make_oxml_element('a:solidFill')
            srgbClr = _make_oxml_element('a:srgbClr')
            srgbClr.set('val', color)
            solidFill.append(srgbClr)
            ln.append(solidFill)

            prstDash = _make_oxml_element('a:prstDash')
            prstDash.set('val', 'solid')
            ln.append(prstDash)

            tcPr.append(ln)

    def _format_table_cell(self, cell, text="", is_header=False, font_size=None,
                           alignment=None, bold=None, text_color=None, bg_color=None):
        """Format a table cell with LG styling."""
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = alignment or PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text

        if is_header:
            self._set_font(run, size=font_size or LGTypography.TABLE_HEADER,
                          bold=True if bold is None else bold,
                          color=text_color or LGColors.WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color or LGColors.CHARCOAL
        else:
            self._set_font(run, size=font_size or LGTypography.TABLE_BODY,
                          bold=bold or False,
                          color=text_color or LGColors.BLACK)
            if bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
            else:
                cell.fill.background()

        # Margins
        cell.margin_left = Cm(0.2)
        cell.margin_right = Cm(0.2)
        cell.margin_top = Cm(0.1)
        cell.margin_bottom = Cm(0.1)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ─────────────────────────────────────────
    # Slide Methods
    # ─────────────────────────────────────────

    def add_cover(self, title, subtitle="", date="", logo_path=None):
        """
        Add a cover slide with L-bracket decorations.

        Args:
            title: Main presentation title
            subtitle: Team/department name
            date: Presentation date (e.g., "2025.10.23")
            logo_path: Override logo path for this slide
        """
        slide = self._get_blank_slide()

        # L-brackets
        self._add_l_bracket(slide, "top-left")
        self._add_l_bracket(slide, "bottom-right")

        # Title (centered)
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        title_box = self._add_textbox(
            slide,
            left=Cm(3), top=sh // 2 - Cm(3),
            width=sw - Cm(6), height=Cm(4),
            text=title,
            size=LGTypography.COVER_TITLE,
            bold=True,
            color=LGColors.BLACK,
            alignment=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE
        )

        # Subtitle + Date (bottom center)
        info_parts = []
        if subtitle:
            info_parts.append(subtitle)
        if date:
            info_parts.append(date)

        if info_parts:
            info_box = self._add_textbox(
                slide,
                left=Cm(3), top=sh - Cm(4),
                width=sw - Cm(6), height=Cm(2.5),
                alignment=PP_ALIGN.CENTER
            )
            tf = info_box.text_frame
            tf.clear()

            if subtitle:
                p1 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p1.alignment = PP_ALIGN.CENTER
                run1 = p1.add_run()
                run1.text = subtitle
                self._set_font(run1, size=LGTypography.COVER_SUBTITLE,
                              color=LGColors.DARK_GRAY)

            if date:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(4)
                run2 = p2.add_run()
                run2.text = date
                self._set_font(run2, size=LGTypography.COVER_SUBTITLE,
                              bold=True, color=LGColors.BLACK)

        # Logo (if provided)
        logo = logo_path or self.logo_path
        if logo and os.path.exists(logo):
            slide.shapes.add_picture(
                logo,
                sw - Cm(5), sh - Cm(4),
                width=Cm(3)
            )

        return slide

    def add_toc(self, items):
        """
        Add a Table of Contents slide.

        Args:
            items: List of tuples (title, [sub_items])
                   e.g., [("Summary", []), ("시스템 소개", ["항목1", "항목2"])]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # Top gray line
        self._add_horizontal_line(
            slide, Cm(1.5), Cm(1.5), sw - Cm(3)
        )

        # "Contents" title
        contents_box = self._add_textbox(
            slide,
            left=Cm(2), top=Cm(2.2),
            width=Cm(10), height=Cm(1.5),
            text="Contents",
            size=LGTypography.TOC_TITLE,
            color=LGColors.BLACK
        )

        # Red underline bar
        red_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(2), Cm(4.0), Cm(2.5), Cm(0.15)
        )
        red_bar.fill.solid()
        red_bar.fill.fore_color.rgb = LGColors.RED
        red_bar.line.fill.background()

        # Gray divider line below title
        self._add_horizontal_line(
            slide, Cm(2), Cm(4.5), sw - Cm(4)
        )

        # TOC items
        y_pos = Cm(5.0)
        roman = ['I', 'II', 'III', 'IV', 'V', 'VI', 'VII', 'VIII', 'IX', 'X']

        for i, (title, sub_items) in enumerate(items):
            # Roman numeral + title
            numeral = roman[i] if i < len(roman) else str(i + 1)
            toc_text = f"{numeral}.   {title}"

            toc_box = self._add_textbox(
                slide,
                left=Cm(3), top=y_pos,
                width=Cm(20), height=Cm(0.9),
                text=toc_text,
                size=LGTypography.TOC_ITEM,
                bold=True,
                color=LGColors.RED
            )
            y_pos += Cm(0.85)

            # Sub-items
            if sub_items:
                for sub in sub_items:
                    sub_box = self._add_textbox(
                        slide,
                        left=Cm(4.5), top=y_pos,
                        width=Cm(20), height=Cm(0.6),
                        text=f"- {sub}",
                        size=LGTypography.TOC_SUBITEM,
                        color=LGColors.DARK_GRAY
                    )
                    y_pos += Cm(0.55)

            y_pos += Cm(0.2)

        return slide

    def add_section_divider(self, number, title, color=None):
        """
        Add a section divider slide.

        Args:
            number: Section number (e.g., "I", "II", "01")
            title: Section title
            color: Override background color
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # Left colored block (1/3 of slide)
        block_width = sw // 3
        block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, block_width, sh
        )
        block.fill.solid()
        block.fill.fore_color.rgb = color or LGColors.RED
        block.line.fill.background()
        self._add_shadow(block)

        # Section number on colored block
        self._add_textbox(
            slide,
            left=Cm(2), top=sh // 2 - Cm(3),
            width=block_width - Cm(3), height=Cm(2),
            text=str(number),
            size=Pt(48),
            bold=True,
            color=LGColors.WHITE,
            alignment=PP_ALIGN.LEFT
        )

        # Title on white area
        self._add_textbox(
            slide,
            left=block_width + Cm(2), top=sh // 2 - Cm(2),
            width=sw - block_width - Cm(4), height=Cm(3),
            text=title,
            size=Pt(36),
            bold=True,
            color=LGColors.BLACK,
            alignment=PP_ALIGN.LEFT,
            vertical=MSO_ANCHOR.MIDDLE
        )

        return slide

    def add_content(self, title, section="", body="", bullets=None,
                    sub_title=""):
        """
        Add a standard content slide with left accent bar.

        Args:
            title: Slide title (e.g., "1.1 시스템 개요")
            section: Section name shown in top-right
            body: Body text paragraph
            bullets: List of bullet point strings
            sub_title: Subtitle text below title
        """
        slide = self._get_blank_slide()

        # Left accent bar
        self._add_accent_bar(slide)

        # Section indicator (top-right)
        if section:
            self._add_section_indicator(slide, section)

        # Title
        self._add_slide_title(slide, title)

        # Subtitle
        y_pos = Cm(1.8)
        if sub_title:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=y_pos,
                width=Cm(28), height=Cm(1.0),
                text=sub_title,
                size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Body text
        if body:
            body_box = self._add_textbox(
                slide,
                left=Cm(1.5), top=y_pos,
                width=Cm(31), height=Cm(12),
                text=body,
                size=Pt(14),
                color=LGColors.BLACK
            )
            y_pos += Cm(2.0)

        # Bullet points
        if bullets:
            # Determine if few bullets → use larger font and spacing
            few_bullets = len(bullets) <= 4 and not body

            # Add thin LG RED divider line between title and bullets for few-bullet slides
            if few_bullets:
                divider = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Cm(1.5), int(y_pos) - int(Cm(0.15)),
                    Cm(31), Pt(1.5)
                )
                divider.fill.solid()
                divider.fill.fore_color.rgb = LGColors.RED
                divider.line.fill.background()
                y_pos = int(y_pos) + int(Cm(0.2))

            bullet_box = slide.shapes.add_textbox(
                Cm(1.5), y_pos, Cm(31), Cm(14)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            bullet_box.fill.background()
            bullet_box.line.fill.background()

            bullet_font_size = Pt(14) if few_bullets else LGTypography.BODY
            bullet_space_after = Pt(8) if few_bullets else Pt(3)

            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.space_after = bullet_space_after

                # Bullet character
                pPr = p._pPr
                if pPr is None:
                    pPr = p._p.get_or_add_pPr()
                buChar = _make_oxml_element('a:buChar')
                buChar.set('char', '\u2022')
                pPr.append(buChar)

                # Indent
                pPr.set('marL', str(int(Cm(0.5))))
                pPr.set('indent', str(int(Cm(-0.4))))

                run = p.add_run()
                run.text = bullet
                self._set_font(run, size=bullet_font_size, color=LGColors.BLACK)

        return slide

    def add_roadmap(self, title, section="", subtitle="", years=None,
                    roadmap_items=None, table_data=None):
        """
        Add a roadmap timeline slide.

        Args:
            title: Slide title (e.g., "[SPC] 로드맵 및 계열사별 비교")
            section: Section name for top-right indicator
            subtitle: Description text below title
            years: List of year/phase strings for timeline headers
                   e.g., ["(2025) 데이터 안정화", "(2026) 분석 고도화", "(2027) 자동화"]
            roadmap_items: Dict with structure:
                {
                    "label": "시스템 로드맵",  # Left label text
                    "rows": [
                        {
                            "items_by_year": [
                                # Items for each year column
                                [{"text": "항목", "tag": "LGES", "tag_color": "#1565C0"}],
                                [{"text": "항목2", "color": "green"}],
                                [{"text": "항목3", "color": "orange"}]
                            ]
                        }
                    ]
                }
            table_data: Dict with structure:
                {
                    "title": "계열사 별 현황",
                    "headers": ["계열사", "DX 수준", "AI 적용", ...],
                    "rows": [["에너지솔루션", "Lv3", "적용(예정)", ...], ...]
                }
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # Left accent bar
        self._add_accent_bar(slide)

        # Section indicator
        if section:
            self._add_section_indicator(slide, section)

        # Title
        self._add_slide_title(slide, title)

        # Subtitle
        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=y_pos,
                width=Cm(28), height=Cm(1.0),
                text=subtitle,
                size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Timeline headers (chevron/arrow style)
        if years:
            num_years = len(years)
            header_left = Cm(4.5)
            total_width = sw - int(header_left) - int(Cm(1.5))
            col_width = total_width // num_years

            for i, year_text in enumerate(years):
                x = int(header_left) + (col_width * i)
                header_bar = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    x, int(y_pos), col_width, Cm(1.2)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = LGColors.CHARCOAL
                header_bar.line.fill.background()
                self._add_shadow(header_bar)

                # Year text inside header
                tf = header_bar.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = year_text
                self._set_font(run, size=Pt(11), bold=True, color=LGColors.WHITE)

            y_pos += Cm(1.8)

        # Roadmap content grid
        if roadmap_items:
            label_text = roadmap_items.get("label", "")
            rows = roadmap_items.get("rows", [])

            # Left label block
            if label_text and years:
                label_height = Cm(len(rows) * 4.5) if rows else Cm(4)
                label_block = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Cm(1.5), int(y_pos), Cm(2.5), label_height
                )
                label_block.fill.solid()
                label_block.fill.fore_color.rgb = LGColors.RED
                label_block.line.fill.background()
                self._add_shadow(label_block)

                # Label text (vertical-ish)
                tf = label_block.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                for char_idx, line in enumerate(label_text.split('\n') if '\n' in label_text else [label_text]):
                    if char_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    self._set_font(run, size=Pt(11), bold=True, color=LGColors.WHITE)

            # Content cells
            if years and rows:
                content_left = Cm(4.5)
                total_width = sw - int(content_left) - int(Cm(1.5))
                col_width = total_width // len(years)

                for row_idx, row_data in enumerate(rows):
                    items_by_year = row_data.get("items_by_year", [])
                    row_top = int(y_pos) + (int(Cm(4.5)) * row_idx)

                    for col_idx, year_items in enumerate(items_by_year):
                        cell_x = int(content_left) + (col_width * col_idx)

                        # Cell background
                        cell_bg = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            cell_x + Cm(0.1), row_top + Cm(0.1),
                            col_width - Cm(0.2), Cm(4.0)
                        )
                        cell_bg.fill.solid()
                        cell_bg.fill.fore_color.rgb = LGColors.LIGHT_GRAY
                        cell_bg.line.fill.background()
                        self._add_shadow(cell_bg)

                        # Cell content
                        if year_items:
                            cell_box = slide.shapes.add_textbox(
                                cell_x + Cm(0.4), row_top + Cm(0.4),
                                col_width - Cm(0.8), Cm(3.5)
                            )
                            cell_tf = cell_box.text_frame
                            cell_tf.word_wrap = True
                            cell_box.fill.background()
                            cell_box.line.fill.background()

                            for item_idx, item in enumerate(year_items):
                                if item_idx == 0:
                                    p = cell_tf.paragraphs[0]
                                else:
                                    p = cell_tf.add_paragraph()

                                p.space_after = Pt(3)

                                # Determine text color
                                item_text = item if isinstance(item, str) else item.get("text", "")
                                item_color_name = None if isinstance(item, str) else item.get("color", None)

                                if item_color_name == "green":
                                    text_color = LGColors.GREEN
                                elif item_color_name == "orange":
                                    text_color = LGColors.ORANGE
                                else:
                                    text_color = LGColors.BLACK

                                # Bullet
                                pPr = p._p.get_or_add_pPr()
                                buChar = _make_oxml_element('a:buChar')
                                buChar.set('char', '\u2022')
                                pPr.append(buChar)

                                run = p.add_run()
                                run.text = item_text
                                self._set_font(run, size=LGTypography.BODY_SMALL,
                                              color=text_color)

                                # Tag badge (e.g., "LGES")
                                if isinstance(item, dict) and item.get("tag"):
                                    tag_run = p.add_run()
                                    tag_run.text = f"  {item['tag']}"
                                    tag_color_hex = item.get("tag_color", "#1565C0")
                                    r, g, b = int(tag_color_hex[1:3], 16), int(tag_color_hex[3:5], 16), int(tag_color_hex[5:7], 16)
                                    self._set_font(tag_run, size=Pt(9), bold=True,
                                                  color=RGBColor(r, g, b))

        # Table section
        if table_data:
            self._add_roadmap_table(slide, table_data, y_start=y_pos)

        return slide

    def _add_roadmap_table(self, slide, table_data, y_start=None):
        """Add a comparison table at the bottom of a slide."""
        title = table_data.get("title", "")
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers or not rows:
            return

        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # Table title
        table_y = sh - Cm(1.5) - Cm(len(rows) * 1.0 + 1.5)
        if title:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=table_y - Cm(1.2),
                width=Cm(15), height=Cm(1.0),
                text=f"\u25a0 {title}",
                size=LGTypography.BODY,
                bold=True,
                color=LGColors.BLACK
            )

        # Table
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        table_width = sw - int(Cm(3))
        table_height = Cm(num_rows * 1.0)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), table_y,
            table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Set column widths evenly
        col_width = table_width // num_cols
        for c in range(num_cols):
            table.columns[c].width = col_width

        # Header row
        for c, header in enumerate(headers):
            self._format_table_cell(table.cell(0, c), header, is_header=True)
            self._set_cell_border(table.cell(0, c), color="3C3C3C")

        # Data rows
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                bg = LGColors.WHITE if r % 2 == 0 else LGColors.LIGHT_GRAY
                self._format_table_cell(table.cell(r + 1, c), str(val), bg_color=bg)
                self._set_cell_border(table.cell(r + 1, c))

    def add_table(self, title, section="", headers=None, rows=None,
                  col_widths=None, subtitle=""):
        """
        Add a dedicated table slide.

        Args:
            title: Slide title
            section: Section name for top-right
            headers: List of column header strings
            rows: List of row data (list of lists)
            col_widths: Optional list of column widths in Cm
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # Left accent bar
        self._add_accent_bar(slide)

        # Section indicator
        if section:
            self._add_section_indicator(slide, section)

        # Title
        self._add_slide_title(slide, title)

        # Subtitle
        y_pos = Cm(2.2)
        if subtitle:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=Cm(1.8),
                width=Cm(28), height=Cm(1.0),
                text=subtitle,
                size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)

        if not headers or not rows:
            return slide

        # Table
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_width = sw - int(Cm(2.0))
        available_height = int(LGDimensions.SLIDE_HEIGHT) - int(y_pos) - int(Cm(0.5))
        row_height = min(Cm(1.2), available_height // num_rows)
        table_height = row_height * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), int(y_pos),
            table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Set column widths
        if col_widths:
            for c, w in enumerate(col_widths):
                if c < num_cols:
                    table.columns[c].width = Cm(w)
        else:
            col_width = table_width // num_cols
            for c in range(num_cols):
                table.columns[c].width = col_width

        # Header row
        for c, header in enumerate(headers):
            self._format_table_cell(table.cell(0, c), header, is_header=True)
            self._set_cell_border(table.cell(0, c), color="3C3C3C")

        # Data rows
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                bg = LGColors.WHITE if r % 2 == 0 else LGColors.LIGHT_GRAY
                self._format_table_cell(table.cell(r + 1, c), str(val), bg_color=bg)
                self._set_cell_border(table.cell(r + 1, c))

        return slide

    def add_blank_content(self, title, section="", subtitle=""):
        """
        Add a blank content slide with just the chrome (accent bar, title, section).
        Returns the slide object for custom content addition.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if subtitle:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=Cm(1.8),
                width=Cm(28), height=Cm(1.0),
                text=subtitle,
                size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY,
                weight="semibold"
            )

        return slide

    def add_recommendation(self, title="제언", section="", subtitle="",
                           recommendations=None):
        """
        Add a recommendation (제언) slide — commonly placed before
        the closing slide. Displays numbered recommendation items
        with optional detail text.

        Args:
            title: Slide title (default: "제언")
            section: Section name for top-right
            subtitle: Optional subtitle
            recommendations: List of strings or list of dicts:
                ["데이터 품질 확보가 최우선", "단계적 MLOps 성숙도 향상"]
                or
                [
                    {"title": "데이터 품질 확보", "detail": "데이터 정합성 검증 체계 구축"},
                    {"title": "MLOps 성숙도 향상", "detail": "단계적 자동화 파이프라인 도입"},
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        recommendations = recommendations or []
        if not recommendations:
            return slide

        # Numbered recommendation items
        item_y = int(y_pos) + int(Cm(0.3))
        for idx, rec in enumerate(recommendations):
            if isinstance(rec, dict):
                rec_title = rec.get("title", "")
                rec_detail = rec.get("detail", "")
            else:
                rec_title = str(rec)
                rec_detail = ""

            num_label = str(idx + 1)

            # Number circle (LG RED)
            circle_size = Cm(1.2)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Cm(1.5), item_y, int(circle_size), int(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = LGColors.RED
            circle.line.fill.background()
            tf = circle.text_frame
            tf.clear()
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num_label
            self._set_font(run, size=Pt(14), bold=True, color=LGColors.WHITE)

            # Recommendation title text
            text_left = Cm(3.2)
            self._add_textbox(
                slide, text_left, item_y,
                Cm(29), Cm(1.2),
                text=rec_title,
                size=Pt(16), bold=True,
                color=LGColors.BLACK,
                alignment=PP_ALIGN.LEFT,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Detail text (if provided)
            if rec_detail:
                self._add_textbox(
                    slide, text_left, item_y + int(Cm(1.2)),
                    Cm(29), Cm(1.0),
                    text=rec_detail,
                    size=LGTypography.BODY,
                    color=LGColors.DARK_GRAY,
                    alignment=PP_ALIGN.LEFT,
                    vertical=MSO_ANCHOR.TOP
                )
                item_y += int(Cm(2.5))
            else:
                item_y += int(Cm(1.8))

        return slide

    def add_closing(self, text="감사합니다", subtitle=""):
        """
        Add a closing/thank you slide.

        Args:
            text: Main closing text
            subtitle: Additional text
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-brackets (same as cover)
        self._add_l_bracket(slide, "top-left")
        self._add_l_bracket(slide, "bottom-right")

        # Closing text (centered)
        self._add_textbox(
            slide,
            left=Cm(3), top=sh // 2 - Cm(2),
            width=sw - Cm(6), height=Cm(3),
            text=text,
            size=Pt(36),
            bold=True,
            color=LGColors.BLACK,
            alignment=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE
        )

        if subtitle:
            self._add_textbox(
                slide,
                left=Cm(3), top=sh // 2 + Cm(1.5),
                width=sw - Cm(6), height=Cm(1.5),
                text=subtitle,
                size=LGTypography.COVER_SUBTITLE,
                color=LGColors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        # Logo
        if self.logo_path and os.path.exists(self.logo_path):
            slide.shapes.add_picture(
                self.logo_path,
                sw - Cm(5), sh - Cm(4),
                width=Cm(3)
            )

        return slide

    # ─────────────────────────────────────────
    # Summary Matrix / Pitch Deck Templates
    # ─────────────────────────────────────────

    def add_summary_matrix(self, title, section="", headers=None,
                           row_groups=None, subtitle=""):
        """
        Add a summary matrix table — grouped rows with merged category cells.
        Matches the LG pattern of category labels on the left with detailed
        content cells across columns.

        Args:
            title: Slide title (e.g., "공정 DX팀 시스템 현황 Summary")
            section: Section name for top-right indicator
            headers: List of column header strings (e.g., ["LG에너지솔루션", "LG디스플레이", ...])
            row_groups: List of dicts, each representing a merged group:
                [
                    {
                        "category": "공정제어",      # Left-most merged label
                        "rows": [
                            {
                                "sub_label": "과제\\n현황",  # Sub-category label
                                "cells": ["셀1 내용", "셀2 내용", ...]  # One per header
                            },
                            {
                                "sub_label": "AI 적용\\n과제",
                                "cells": ["셀1 내용", "셀2 내용", ...]
                            }
                        ]
                    }
                ]
            subtitle: Optional subtitle text
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # Accent block + title
        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if not headers or not row_groups:
            return slide

        # Calculate table dimensions
        # Columns: category_col + sub_label_col + data_cols
        num_data_cols = len(headers)
        num_cols = 2 + num_data_cols  # category + sub_label + data columns
        total_data_rows = sum(len(g["rows"]) for g in row_groups)
        num_rows = 1 + total_data_rows  # header + data

        y_start = Cm(2.2) if not subtitle else Cm(2.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )

        table_width = sw - int(Cm(2.0))
        available_h = int(LGDimensions.SLIDE_HEIGHT) - int(y_start) - int(Cm(0.8))
        row_h = min(Cm(2.5), available_h // num_rows)
        table_height = row_h * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), int(y_start), table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Column widths
        cat_width = int(table_width * 0.07)
        sub_width = int(table_width * 0.07)
        data_width = (table_width - cat_width - sub_width) // num_data_cols
        table.columns[0].width = cat_width
        table.columns[1].width = sub_width
        for c in range(num_data_cols):
            table.columns[2 + c].width = data_width

        # Header row: empty for first 2 cols, then data headers
        self._format_table_cell(table.cell(0, 0), "", is_header=False,
                                bg_color=LGColors.WHITE)
        self._set_cell_border(table.cell(0, 0), color="CCCCCC")
        self._format_table_cell(table.cell(0, 1), "", is_header=False,
                                bg_color=LGColors.WHITE)
        self._set_cell_border(table.cell(0, 1), color="CCCCCC")
        for c, header in enumerate(headers):
            self._format_table_cell(table.cell(0, 2 + c), header, is_header=True)
            self._set_cell_border(table.cell(0, 2 + c), color="3C3C3C")

        # Data rows with merged category cells
        current_row = 1
        for group in row_groups:
            category = group["category"]
            group_rows = group["rows"]
            group_start = current_row
            group_end = current_row + len(group_rows) - 1

            for i, row_data in enumerate(group_rows):
                sub_label = row_data.get("sub_label", "")
                cells = row_data.get("cells", [])
                r = current_row

                # Sub-label cell (gray background)
                self._format_table_cell(
                    table.cell(r, 1), sub_label.replace("\\n", "\n"),
                    font_size=Pt(9), bold=True, bg_color=LGColors.LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER
                )
                self._set_cell_border(table.cell(r, 1), color="CCCCCC")

                # Data cells
                for c, cell_text in enumerate(cells):
                    if c < num_data_cols:
                        self._format_table_cell(
                            table.cell(r, 2 + c), cell_text,
                            font_size=Pt(9), alignment=PP_ALIGN.LEFT,
                            bg_color=LGColors.WHITE
                        )
                        self._set_cell_border(table.cell(r, 2 + c), color="CCCCCC")

                current_row += 1

            # Merge category cells vertically
            if len(group_rows) > 1:
                table.cell(group_start, 0).merge(table.cell(group_end, 0))
            self._format_table_cell(
                table.cell(group_start, 0), category,
                font_size=Pt(10), bold=True, bg_color=LGColors.LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER
            )
            self._set_cell_border(table.cell(group_start, 0), color="CCCCCC")

        return slide

    def add_two_column(self, title, section="", left_title="", left_bullets=None,
                       right_title="", right_bullets=None, subtitle=""):
        """
        Add a two-column layout slide — useful for comparisons, pros/cons,
        before/after, or side-by-side content.

        Args:
            title: Slide title
            section: Section name for top-right
            left_title: Title for left column
            left_bullets: List of bullet strings for left column
            right_title: Title for right column
            right_bullets: List of bullet strings for right column
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(2.2)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)

        col_width = (sw - int(Cm(3.5))) // 2
        left_x = Cm(1.5)
        right_x = int(left_x) + col_width + int(Cm(0.5))
        header_h = Cm(1.0)
        body_h = int(LGDimensions.SLIDE_HEIGHT) - int(y_pos) - int(header_h) - int(Cm(1.5))

        for col_x, col_title, col_bullets in [
            (left_x, left_title, left_bullets or []),
            (right_x, right_title, right_bullets or [])
        ]:
            # Column header box (CHARCOAL background with white text)
            if col_title:
                hdr_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(col_x), int(y_pos), col_width, int(header_h)
                )
                hdr_box.fill.solid()
                hdr_box.fill.fore_color.rgb = LGColors.CHARCOAL
                hdr_box.line.fill.background()
                tf = hdr_box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Cm(0.4)
                tf.margin_right = Cm(0.3)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = col_title
                self._set_font(run, size=Pt(14), bold=True,
                              color=LGColors.WHITE)

            # Column body box (Light gray background with bullets)
            body_top = int(y_pos) + int(header_h)
            body_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(col_x), body_top, col_width, body_h
            )
            body_box.fill.solid()
            body_box.fill.fore_color.rgb = LGColors.LIGHT_GRAY
            body_box.line.fill.background()

            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.4)
            tf.margin_right = Cm(0.3)
            tf.margin_top = Cm(0.3)

            if col_bullets:
                for j, bullet in enumerate(col_bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '\u2022')
                    pPr.append(buChar)
                    pPr.set('marL', str(int(Cm(0.5))))
                    pPr.set('indent', str(int(Cm(-0.4))))
                    run = p.add_run()
                    run.text = bullet
                    self._set_font(run, size=LGTypography.BODY,
                                  color=LGColors.BLACK)

            self._add_shadow(body_box)

        return slide

    def add_kpi_cards(self, title, section="", kpis=None, subtitle=""):
        """
        Add a KPI/metrics card slide — large numbers with labels,
        commonly used in executive summaries and dashboards.

        Args:
            title: Slide title
            section: Section name for top-right
            kpis: List of dicts with KPI data:
                [
                    {"value": "30%", "label": "불량률 감소", "color": "#2E7D32"},
                    {"value": "15%", "label": "공정 효율 향상"},
                    {"value": "2.5억", "label": "연간 비용 절감", "color": "#1565C0"},
                ]
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if not kpis:
            return slide

        y_pos = Cm(3.0)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(3.5)

        num_cards = len(kpis)
        card_gap = Cm(0.3)
        total_gap = int(card_gap) * (num_cards - 1)
        card_width = (sw - int(Cm(2.0)) - total_gap) // num_cards
        card_height = Cm(8)
        x_start = Cm(1.5)

        for i, kpi in enumerate(kpis):
            x = int(x_start) + i * (card_width + int(card_gap))

            # Card background
            card = self.add_box(
                slide, x, int(y_pos), card_width, card_height,
                bg_color=LGColors.LIGHT_GRAY
            )

            # Value (large number)
            value_color_hex = kpi.get("color", None)
            if value_color_hex:
                r, g, b = int(value_color_hex[1:3], 16), int(value_color_hex[3:5], 16), int(value_color_hex[5:7], 16)
                v_color = RGBColor(r, g, b)
            else:
                v_color = LGColors.RED

            self._add_textbox(
                slide, x + Cm(0.5), int(y_pos) + Cm(1.5),
                card_width - Cm(1.0), Cm(3.5),
                text=kpi.get("value", ""),
                size=Pt(40), bold=True, color=v_color,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Label
            self._add_textbox(
                slide, x + Cm(0.5), int(y_pos) + Cm(5.5),
                card_width - Cm(1.0), Cm(2.0),
                text=kpi.get("label", ""),
                size=LGTypography.BODY, color=LGColors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.TOP
            )

        return slide

    def add_timeline(self, title, section="", milestones=None, subtitle=""):
        """
        Add a horizontal timeline slide — useful for project phases,
        history, or sequential milestones.

        Args:
            title: Slide title
            section: Section name for top-right
            milestones: List of dicts:
                [
                    {"date": "2025 Q1", "title": "Phase 1", "description": "설명"},
                    {"date": "2025 Q2", "title": "Phase 2", "description": "설명"},
                ]
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if not milestones:
            return slide

        y_pos = Cm(2.2)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)

        num = len(milestones)
        line_y = int(y_pos) + int(Cm(4.0))
        margin_x = Cm(1.5)
        line_width = sw - int(margin_x) * 2

        # Horizontal line
        self._add_horizontal_line(
            slide, margin_x, line_y, line_width,
            color=LGColors.CHARCOAL, height=Pt(3)
        )

        # Milestones
        spacing = line_width // num
        for i, ms in enumerate(milestones):
            cx = int(margin_x) + spacing * i + spacing // 2

            # Circle marker
            dot_size = Cm(0.6)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx - int(dot_size) // 2, line_y - int(dot_size) // 2,
                dot_size, dot_size
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = LGColors.RED
            dot.line.fill.background()

            # Date above line
            self._add_textbox(
                slide, cx - spacing // 2, line_y - Cm(2.5),
                spacing, Cm(2.0),
                text=ms.get("date", ""),
                size=Pt(11), bold=True, color=LGColors.RED,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.BOTTOM
            )

            # Title + description below line
            self._add_textbox(
                slide, cx - spacing // 2, line_y + Cm(1.0),
                spacing, Cm(1.5),
                text=ms.get("title", ""),
                size=LGTypography.BODY, bold=True, color=LGColors.BLACK,
                alignment=PP_ALIGN.CENTER
            )

            if ms.get("description"):
                self._add_textbox(
                    slide, cx - spacing // 2, line_y + Cm(2.5),
                    spacing, Cm(4.0),
                    text=ms["description"],
                    size=LGTypography.BODY_SMALL, color=LGColors.DARK_GRAY,
                    alignment=PP_ALIGN.CENTER
                )

        return slide

    def add_process_flow(self, title, section="", steps=None, subtitle=""):
        """
        Add a process flow slide — horizontal arrow-connected steps.
        Good for workflows, system architecture overview, or methodology.

        Args:
            title: Slide title
            section: Section name for top-right
            steps: List of dicts:
                [
                    {"title": "데이터 수집", "items": ["센서 데이터", "로그 수집"]},
                    {"title": "전처리", "items": ["정규화", "이상치 제거"]},
                    {"title": "AI 분석", "items": ["모델 학습", "추론"]},
                    {"title": "대시보드", "items": ["시각화", "알림"]},
                ]
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if not steps:
            return slide

        y_pos = Cm(2.2)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)

        num = len(steps)
        arrow_width = Cm(0.8)

        # Serpentine layout for >4 steps
        if num > 4:
            row1_count = (num + 1) // 2  # first row gets half (rounded up)
            row2_count = num - row1_count
            total_arrows_r1 = int(arrow_width) * (row1_count - 1)
            step_width = (sw - int(Cm(3.0)) - total_arrows_r1) // row1_count
            step_height = Cm(4.5)
            row_gap = Cm(1.2)
            x_start = Cm(1.5)

            # --- Row 1: left to right ---
            for i in range(row1_count):
                step = steps[i]
                box_x = int(x_start) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, int(y_pos), step_width, Cm(1.2),
                    text=step.get("title", ""),
                    bg_color=LGColors.CHARCOAL, text_color=LGColors.WHITE,
                    font_size=LGTypography.BODY, bold=True
                )
                content = self.add_box(
                    slide, box_x, int(y_pos) + int(Cm(1.2)),
                    step_width, step_height - Cm(1.2),
                    bg_color=LGColors.LIGHT_GRAY
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), int(y_pos) + int(Cm(1.6)),
                        step_width - Cm(0.6), step_height - Cm(2.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=LGColors.BLACK)

                if i < row1_count - 1:
                    arrow_x = box_x + step_width
                    arrow_y = int(y_pos) + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = LGColors.CHARCOAL
                    arrow.line.fill.background()

            # Down arrow from last in row1 to first in row2
            last_r1_x = int(x_start) + (row1_count - 1) * (step_width + int(arrow_width))
            down_arrow_x = last_r1_x + step_width // 2 - int(Cm(0.5))
            down_arrow_y = int(y_pos) + int(step_height)
            down_arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                down_arrow_x, down_arrow_y, Cm(1.0), row_gap
            )
            down_arrow.fill.solid()
            down_arrow.fill.fore_color.rgb = LGColors.CHARCOAL
            down_arrow.line.fill.background()

            # --- Row 2: left to right ---
            y_pos2 = int(y_pos) + int(step_height) + int(row_gap)
            for i in range(row2_count):
                step = steps[row1_count + i]
                box_x = int(x_start) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, y_pos2, step_width, Cm(1.2),
                    text=step.get("title", ""),
                    bg_color=LGColors.CHARCOAL, text_color=LGColors.WHITE,
                    font_size=LGTypography.BODY, bold=True
                )
                content = self.add_box(
                    slide, box_x, y_pos2 + int(Cm(1.2)),
                    step_width, step_height - Cm(1.2),
                    bg_color=LGColors.LIGHT_GRAY
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), y_pos2 + int(Cm(1.6)),
                        step_width - Cm(0.6), step_height - Cm(2.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=LGColors.BLACK)

                if i < row2_count - 1:
                    arrow_x = box_x + step_width
                    arrow_y = y_pos2 + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = LGColors.CHARCOAL
                    arrow.line.fill.background()

        else:
            # Original layout for <=4 steps
            total_arrows = int(arrow_width) * (num - 1)
            step_width = (sw - int(Cm(3.0)) - total_arrows) // num
            step_height = Cm(8)
            x = Cm(1.5)

            for i, step in enumerate(steps):
                box_x = int(x) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, int(y_pos), step_width, Cm(1.5),
                    text=step.get("title", ""),
                    bg_color=LGColors.CHARCOAL, text_color=LGColors.WHITE,
                    font_size=LGTypography.BODY, bold=True
                )

                content = self.add_box(
                    slide, box_x, int(y_pos) + int(Cm(1.5)),
                    step_width, step_height - Cm(1.5),
                    bg_color=LGColors.LIGHT_GRAY
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), int(y_pos) + int(Cm(2.2)),
                        step_width - Cm(0.6), step_height - Cm(3.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()

                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=LGColors.BLACK)

                if i < num - 1:
                    arrow_x = box_x + step_width
                    arrow_y = int(y_pos) + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = LGColors.CHARCOAL
                    arrow.line.fill.background()

        return slide

    def add_swot(self, title, section="", strengths=None, weaknesses=None,
                 opportunities=None, threats=None, subtitle=""):
        """
        Add a SWOT analysis slide — 2x2 grid with color-coded quadrants.

        Args:
            title: Slide title
            section: Section name for top-right
            strengths: List of strength strings
            weaknesses: List of weakness strings
            opportunities: List of opportunity strings
            threats: List of threat strings
            subtitle: Optional subtitle
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(2.2)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), Cm(1.8), Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE, color=LGColors.DARK_GRAY,
                weight="semibold"
            )
            y_pos = Cm(2.8)

        gap = Cm(0.3)
        quad_w = (sw - int(Cm(3.0)) - int(gap)) // 2
        quad_h = (int(LGDimensions.SLIDE_HEIGHT) - int(y_pos) - int(Cm(0.5)) - int(gap)) // 2
        left_x = Cm(1.5)
        right_x = int(left_x) + quad_w + int(gap)
        top_y = int(y_pos)
        bottom_y = top_y + quad_h + int(gap)

        quadrants = [
            ("Strengths (강점)", strengths or [], left_x, top_y, LGColors.RED),
            ("Weaknesses (약점)", weaknesses or [], right_x, top_y, LGColors.CHARCOAL),
            ("Opportunities (기회)", opportunities or [], left_x, bottom_y, LGColors.GREEN),
            ("Threats (위협)", threats or [], right_x, bottom_y, LGColors.ORANGE),
        ]

        for q_title, items, qx, qy, color in quadrants:
            # Header bar
            self.add_box(
                slide, qx, qy, quad_w, Cm(1.3),
                text=q_title, bg_color=color, text_color=LGColors.WHITE,
                font_size=LGTypography.BODY, bold=True
            )

            # Content area
            self.add_box(
                slide, qx, qy + int(Cm(1.3)), quad_w, quad_h - int(Cm(1.3)),
                bg_color=LGColors.LIGHT_GRAY
            )

            # Items
            if items:
                items_box = slide.shapes.add_textbox(
                    qx + Cm(0.5), qy + int(Cm(1.8)),
                    quad_w - Cm(1.0), quad_h - int(Cm(2.3))
                )
                tf = items_box.text_frame
                tf.word_wrap = True
                items_box.fill.background()
                items_box.line.fill.background()

                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '\u2022')
                    pPr.append(buChar)
                    run = p.add_run()
                    run.text = item
                    self._set_font(run, size=LGTypography.BODY_SMALL,
                                  color=LGColors.BLACK)

        return slide

    # ─────────────────────────────────────────
    # Utility Methods
    # ─────────────────────────────────────────

    def add_box(self, slide, left, top, width, height, text="",
                bg_color=None, text_color=None, font_size=None,
                bold=False, alignment=None, border_color=None,
                shadow=True):
        """
        Add a styled content box to a slide with optional drop shadow.
        Useful for building custom diagram/architecture slides.

        Args:
            shadow: Whether to add a subtle drop shadow (default True).
                    Set to False for flat elements like dividers.

        Returns the shape object.
        """
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color or LGColors.LIGHT_GRAY
        if border_color:
            box.line.color.rgb = border_color
            box.line.width = Pt(1)
        else:
            box.line.fill.background()

        if text:
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.15)
            tf.margin_right = Cm(0.15)
            tf.margin_top = Cm(0.05)
            tf.margin_bottom = Cm(0.05)

            p = tf.paragraphs[0]
            p.alignment = alignment or PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run,
                          size=font_size or LGTypography.BODY,
                          bold=bold,
                          color=text_color or LGColors.BLACK)

        # Add shadow for visual depth (matching LG sample style)
        if shadow:
            self._add_shadow(box)

        return box

    def add_label_badge(self, slide, left, top, text, width=None,
                        bg_color=None, text_color=None):
        """
        Add a small label badge (category tag).
        e.g., red background with white text for section labels.
        """
        w = width or Cm(4)
        badge = self.add_box(
            slide, left, top, w, Cm(0.8),
            text=text,
            bg_color=bg_color or LGColors.RED,
            text_color=text_color or LGColors.WHITE,
            font_size=Pt(9),
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        return badge

    # ─────────────────────────────────────────
    # Advanced Templates (전략보고서/사업계획서)
    # ─────────────────────────────────────────

    def add_architecture(self, title, section="", subtitle="",
                         columns=None, rows=None):
        """
        Add a multi-column architecture/system diagram slide — complex grid
        layout with header row, row groups, and area descriptions.

        Designed for system architecture overviews like the GTM Assistant
        structure or similar multi-layer diagrams seen in strategy decks.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Description text (can be multi-line)
            columns: List of column header strings,
                     e.g. ["MCP Server", "주요 MCP Tool", "Sources/API", "Area"]
            rows: List of row dicts:
                [
                    {
                        "label": "Process User Input",  # left label
                        "cells": [
                            "제품/서비스정보 Read\\nInput Validation",
                            "",  # empty = skip
                            "Percept\\nAgentic AI가 목표 달성을 위해\\n제품/서비스 정보를 인식하고 이해하는 영역"
                        ],
                        "highlight": False  # optional: highlight this row group
                    },
                    ...
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            # Support multi-line subtitle
            sub_box = slide.shapes.add_textbox(
                Cm(1.5), y_pos, Cm(31), Cm(1.5)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            sub_box.fill.background()
            sub_box.line.fill.background()
            for i, line in enumerate(subtitle.split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                self._set_font(run, size=LGTypography.SUBTITLE,
                              color=LGColors.DARK_GRAY, weight="semibold")
            y_pos += Cm(1.8)
        else:
            y_pos = Cm(2.2)

        columns = columns or []
        rows = rows or []
        if not columns:
            return slide

        num_cols = len(columns)
        grid_left = Cm(1.5)
        grid_width = sw - int(Cm(2.0))

        # Column widths: first column slightly narrower for labels
        first_col_w = int(grid_width * 0.18)
        remaining_w = grid_width - first_col_w
        col_w = remaining_w // (num_cols - 1) if num_cols > 1 else remaining_w

        # Header row
        header_h = Cm(1.0)
        x = int(grid_left)
        for ci, col_name in enumerate(columns):
            w = first_col_w if ci == 0 else col_w
            self.add_box(slide, x, int(y_pos), w, int(header_h),
                        text=col_name,
                        bg_color=LGColors.CHARCOAL,
                        text_color=LGColors.WHITE,
                        font_size=Pt(10), bold=True, shadow=False)
            x += w
        y_pos += Cm(1.1)

        # Content rows
        for row in rows:
            label = row.get("label", "")
            cells = row.get("cells", [])
            highlight = row.get("highlight", False)

            # Ensure cells count matches data columns (pad/truncate)
            expected_cells = num_cols - 1
            cells = (cells + [''] * expected_cells)[:expected_cells]

            # Calculate row height based on content
            max_lines = 1
            for cell_text in cells:
                if cell_text:
                    lines = cell_text.count('\n') + 1
                    max_lines = max(max_lines, lines)
            row_h = max(int(Cm(1.2)), int(Cm(0.5 + max_lines * 0.5)))

            x = int(grid_left)
            # First column: label
            label_bg = LGColors.RED if highlight else LGColors.LIGHT_GRAY
            label_tc = LGColors.WHITE if highlight else LGColors.BLACK
            self.add_box(slide, x, int(y_pos), first_col_w, row_h,
                        text=label,
                        bg_color=label_bg, text_color=label_tc,
                        font_size=Pt(9), bold=True, shadow=False)
            x += first_col_w

            # Remaining columns
            for ci, cell_text in enumerate(cells):
                w = col_w
                if cell_text:
                    # Create box with multi-line text
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x, int(y_pos), w, row_h
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = LGColors.LIGHT_GRAY
                    box.line.color.rgb = LGColors.BORDER_GRAY
                    box.line.width = Pt(0.5)

                    tf = box.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for li, line in enumerate(cell_text.split('\n')):
                        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                        p.alignment = PP_ALIGN.CENTER
                        run = p.add_run()
                        run.text = line
                        # First line bold if multiple lines (acts as sub-header)
                        is_bold = (li == 0 and cell_text.count('\n') > 0)
                        self._set_font(run, size=Pt(9),
                                      bold=is_bold,
                                      color=LGColors.BLACK)

                else:
                    # Empty cell - light border only
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x, int(y_pos), w, row_h
                    )
                    box.fill.background()
                    box.line.color.rgb = LGColors.BORDER_GRAY
                    box.line.width = Pt(0.5)
                x += w

            y_pos += row_h + Cm(0.1)

        return slide

    def add_strategy_pillars(self, title, section="", subtitle="",
                             pillars=None):
        """
        Add a strategy pillars slide — 3 to 5 vertical columns showing
        strategic focus areas, commonly used for strategic direction,
        capability frameworks, or organizational pillars.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Description text
            pillars: List of dicts:
                [
                    {
                        "header": "디지털 전환",
                        "header_color": "#A50034",  # optional, default LG RED
                        "items": ["AI/ML 기반 공정 최적화", "디지털 트윈 구축", ...]
                    },
                    ...
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        pillars = pillars or []
        if not pillars:
            return slide

        num = len(pillars)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        pillar_w = (total_w - int(gap) * (num - 1)) // num
        header_h = Cm(1.2)
        body_h = int(LGDimensions.SLIDE_HEIGHT) - int(y_pos) - int(header_h) - int(Cm(0.5))

        x = Cm(1.5)
        for pillar in pillars:
            header = pillar.get("header", "") or pillar.get("title", "")
            items = pillar.get("items", [])
            hdr_color_hex = pillar.get("header_color", None)

            if hdr_color_hex:
                hdr_color = RGBColor(
                    int(hdr_color_hex[1:3], 16),
                    int(hdr_color_hex[3:5], 16),
                    int(hdr_color_hex[5:7], 16)
                )
            else:
                hdr_color = LGColors.CHARCOAL

            # Header
            self.add_box(slide, int(x), int(y_pos), pillar_w, int(header_h),
                        text=header, bg_color=hdr_color,
                        text_color=LGColors.WHITE,
                        font_size=Pt(14), bold=True, shadow=False)

            # Body
            body_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(x), int(y_pos) + int(header_h), pillar_w, body_h
            )
            body_box.fill.solid()
            body_box.fill.fore_color.rgb = LGColors.LIGHT_GRAY
            body_box.line.fill.background()

            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.3)
            tf.margin_right = Cm(0.3)
            tf.margin_top = Cm(0.3)

            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(3)
                # Bullet
                pPr = p._p.get_or_add_pPr()
                buChar = _make_oxml_element('a:buChar')
                buChar.set('char', '\u2022')
                pPr.append(buChar)
                pPr.set('marL', str(int(Cm(0.5))))
                pPr.set('indent', str(int(Cm(-0.4))))

                run = p.add_run()
                run.text = item
                self._set_font(run, size=LGTypography.BODY,
                              color=LGColors.BLACK)

            self._add_shadow(body_box)
            x = int(x) + pillar_w + int(gap)

        return slide

    def add_risk_matrix(self, title, section="", subtitle="",
                        risks=None, x_label="영향도 (Impact)",
                        y_label="발생 가능성 (Likelihood)"):
        """
        Add a risk assessment matrix slide — 3x3 grid with color-coded risk
        levels, commonly used in strategic planning and project risk analysis.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            risks: List of risk items to place on the matrix:
                [
                    {"name": "데이터 유출", "likelihood": 2, "impact": 3},
                    {"name": "인력 이탈", "likelihood": 1, "impact": 2},
                    ...
                ]
                where likelihood and impact are 1 (Low), 2 (Medium), 3 (High)
            x_label: Label for the X axis (impact)
            y_label: Label for the Y axis (likelihood)
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        risks = risks or []

        # Grid layout
        grid_left = Cm(4.5)
        grid_top = int(y_pos) + int(Cm(0.5))
        cell_size = Cm(4.0)
        grid_gap = Cm(0.15)

        # Color coding: [row][col] → (likelihood high→low, impact low→high)
        colors_3x3 = [
            # impact Low, Med, High  (likelihood High = top row)
            [RGBColor(0xD4, 0x76, 0x0A), RGBColor(0xA5, 0x00, 0x34), RGBColor(0xA5, 0x00, 0x34)],
            [RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xD4, 0x76, 0x0A), RGBColor(0xA5, 0x00, 0x34)],
            [RGBColor(0x2E, 0x7D, 0x32), RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xD4, 0x76, 0x0A)],
        ]
        labels_3x3 = [
            ["중간", "높음", "매우 높음"],
            ["낮음", "중간", "높음"],
            ["매우 낮음", "낮음", "중간"],
        ]

        # Y axis label
        self._add_textbox(
            slide, Cm(1.5), grid_top, Cm(2.5), int(cell_size) * 3 + int(grid_gap) * 2,
            text=y_label, size=Pt(10), bold=True, color=LGColors.DARK_GRAY
        )

        # Y axis level labels
        y_levels = ["High", "Medium", "Low"]
        for r in range(3):
            cy = grid_top + r * (int(cell_size) + int(grid_gap))
            self._add_textbox(
                slide, Cm(2.8), cy, Cm(1.5), int(cell_size),
                text=y_levels[r], size=Pt(9), color=LGColors.MEDIUM_GRAY,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # X axis label
        self._add_textbox(
            slide, int(grid_left), grid_top + 3 * (int(cell_size) + int(grid_gap)) + int(Cm(0.3)),
            int(cell_size) * 3 + int(grid_gap) * 2, Cm(1.0),
            text=x_label, size=Pt(10), bold=True, color=LGColors.DARK_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # X axis level labels
        x_levels = ["Low", "Medium", "High"]
        for c in range(3):
            cx = int(grid_left) + c * (int(cell_size) + int(grid_gap))
            self._add_textbox(
                slide, cx,
                grid_top + 3 * (int(cell_size) + int(grid_gap)),
                int(cell_size), Cm(0.8),
                text=x_levels[c], size=Pt(9), color=LGColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        # Draw grid cells
        for r in range(3):
            for c in range(3):
                cx = int(grid_left) + c * (int(cell_size) + int(grid_gap))
                cy = grid_top + r * (int(cell_size) + int(grid_gap))
                cell_color = colors_3x3[r][c]

                box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, cx, cy, int(cell_size), int(cell_size)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = cell_color
                box.line.fill.background()

                # Muted label inside
                tf = box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = labels_3x3[r][c]
                self._set_font(run, size=Pt(8), color=LGColors.WHITE)

        # Place risk items
        for risk in risks:
            name = risk.get("name", "")
            lk = risk.get("likelihood", 1)  # 1-3
            im = risk.get("impact", 1)       # 1-3
            # Map to grid: likelihood 3=top(row0), impact 1=left(col0)
            r = 3 - lk
            c = im - 1
            cx = int(grid_left) + c * (int(cell_size) + int(grid_gap)) + int(Cm(0.3))
            cy = grid_top + r * (int(cell_size) + int(grid_gap)) + int(Cm(0.3))

            # Risk item tag
            self.add_box(
                slide, cx, cy,
                int(cell_size) - int(Cm(0.6)), Cm(1.0),
                text=name, bg_color=LGColors.WHITE,
                text_color=LGColors.BLACK,
                font_size=Pt(8), bold=True,
                shadow=True
            )

        # Legend on the right side
        legend_x = int(grid_left) + 3 * (int(cell_size) + int(grid_gap)) + int(Cm(1.5))
        legend_y = grid_top
        self._add_textbox(
            slide, legend_x, legend_y, Cm(5), Cm(0.8),
            text="위험 수준", size=Pt(10), bold=True, color=LGColors.BLACK
        )
        legend_items = [
            (LGColors.RED, "높음 (High Risk)"),
            (RGBColor(0xD4, 0x76, 0x0A), "중간 (Medium Risk)"),
            (LGColors.GREEN, "낮음 (Low Risk)"),
        ]
        for li, (color, label) in enumerate(legend_items):
            ly = legend_y + int(Cm(1.0)) + li * int(Cm(1.0))
            # Color swatch
            swatch = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, legend_x, ly, Cm(0.6), Cm(0.6)
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = color
            swatch.line.fill.background()
            # Label
            self._add_textbox(
                slide, legend_x + int(Cm(1.0)), ly, Cm(5), Cm(0.6),
                text=label, size=Pt(9), color=LGColors.BLACK
            )

        return slide

    def add_financial_summary(self, title, section="", subtitle="",
                              categories=None, total_label="합계"):
        """
        Add a financial/budget summary slide — structured table with
        category groups, subtotals and a grand total row.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            categories: List of category dicts:
                [
                    {
                        "name": "인건비",
                        "items": [
                            {"name": "개발인력", "2025": "5.0", "2026": "8.0", "2027": "10.0"},
                            {"name": "운영인력", "2025": "2.0", "2026": "3.0", "2027": "4.0"},
                        ],
                        "subtotal": {"2025": "7.0", "2026": "11.0", "2027": "14.0"}
                    },
                    ...
                ]
            total_label: Label for grand total row
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        categories = categories or []
        if not categories:
            return slide

        # Determine year columns from first category
        sample_item = categories[0]["items"][0] if categories[0].get("items") else {}
        year_keys = [k for k in sample_item.keys() if k != "name"]

        # Build table data
        headers = ["구분", "항목"] + [f"{y}년" if len(y) == 4 else y for y in year_keys]
        num_cols = len(headers)

        # Count total rows
        total_rows = 1  # header
        for cat in categories:
            total_rows += len(cat.get("items", []))
            if cat.get("subtotal"):
                total_rows += 1  # subtotal row
        total_rows += 1  # grand total

        table_left = Cm(1.5)
        table_width = sw - int(Cm(2.0))
        row_h = Cm(0.9)

        # Calculate column widths first, then set table width to exact sum
        cat_col_w = int(table_width * 0.12)
        item_col_w = int(table_width * 0.22)
        year_col_w = (table_width - cat_col_w - item_col_w) // max(len(year_keys), 1)
        exact_width = cat_col_w + item_col_w + year_col_w * len(year_keys)

        table_shape = slide.shapes.add_table(
            total_rows, num_cols,
            int(table_left), int(y_pos),
            exact_width, int(row_h) * total_rows
        )
        table = table_shape.table

        # Clear default table style to prevent rendering conflicts
        tbl = table._tbl
        tblPr = tbl.tblPr
        if tblPr is not None:
            tblPr.set('bandRow', '0')
            tblPr.set('bandCol', '0')
            tblPr.set('firstRow', '0')
            tblPr.set('lastRow', '0')
            tblPr.set('firstCol', '0')
            tblPr.set('lastCol', '0')

        # Column widths
        table.columns[0].width = cat_col_w
        table.columns[1].width = item_col_w
        for ci in range(2, num_cols):
            table.columns[ci].width = year_col_w

        # Header row
        for ci, header in enumerate(headers):
            self._format_table_cell(table.cell(0, ci), text=header,
                                   is_header=True)
            self._set_cell_border(table.cell(0, ci))

        # Data rows
        row_idx = 1
        for cat in categories:
            cat_name = cat.get("name", "")
            items = cat.get("items", [])
            subtotal = cat.get("subtotal", None)

            # Category items
            for ii, item in enumerate(items):
                # Category column (merge visual: only first row shows name)
                if ii == 0:
                    self._format_table_cell(
                        table.cell(row_idx, 0), text=cat_name,
                        bold=True, bg_color=LGColors.LIGHT_GRAY,
                        alignment=PP_ALIGN.CENTER
                    )
                else:
                    self._format_table_cell(
                        table.cell(row_idx, 0), text="",
                        bg_color=LGColors.LIGHT_GRAY
                    )

                # Item name
                self._format_table_cell(
                    table.cell(row_idx, 1), text=item.get("name", ""),
                    alignment=PP_ALIGN.LEFT
                )

                # Year values
                for yi, yk in enumerate(year_keys):
                    self._format_table_cell(
                        table.cell(row_idx, 2 + yi),
                        text=str(item.get(yk, "")),
                        alignment=PP_ALIGN.RIGHT
                    )

                for ci in range(num_cols):
                    self._set_cell_border(table.cell(row_idx, ci))
                row_idx += 1

            # Subtotal row
            if subtotal:
                self._format_table_cell(
                    table.cell(row_idx, 0), text="",
                    bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                )
                self._format_table_cell(
                    table.cell(row_idx, 1), text="소계",
                    bold=True, alignment=PP_ALIGN.CENTER,
                    bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                )
                for yi, yk in enumerate(year_keys):
                    self._format_table_cell(
                        table.cell(row_idx, 2 + yi),
                        text=str(subtotal.get(yk, "")),
                        bold=True, alignment=PP_ALIGN.RIGHT,
                        bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                    )
                for ci in range(num_cols):
                    self._set_cell_border(table.cell(row_idx, ci))
                row_idx += 1

        # Grand total row
        self._format_table_cell(
            table.cell(row_idx, 0), text="",
            is_header=True, bg_color=LGColors.CHARCOAL
        )
        self._format_table_cell(
            table.cell(row_idx, 1), text=total_label,
            is_header=True, bg_color=LGColors.CHARCOAL,
            alignment=PP_ALIGN.CENTER
        )
        # Sum subtotals for grand total
        for yi, yk in enumerate(year_keys):
            grand = 0
            for cat in categories:
                st = cat.get("subtotal", {})
                try:
                    grand += float(st.get(yk, 0))
                except (ValueError, TypeError):
                    pass
            self._format_table_cell(
                table.cell(row_idx, 2 + yi),
                text=f"{grand:.1f}",
                is_header=True, bg_color=LGColors.CHARCOAL,
                alignment=PP_ALIGN.RIGHT
            )
        for ci in range(num_cols):
            self._set_cell_border(table.cell(row_idx, ci), color="3C3C3C")

        # Merge category cells vertically
        for cat in categories:
            pass  # Visual merging done by leaving subsequent rows empty

        return slide

    def add_milestone_tracker(self, title, section="", subtitle="",
                              phases=None):
        """
        Add a milestone/progress tracker slide — horizontal chevron phases
        with detailed milestones underneath, useful for project status or
        implementation plans.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            phases: List of phase dicts:
                [
                    {
                        "name": "Phase 1: 준비",
                        "period": "2025.Q1-Q2",
                        "status": "completed",  # "completed", "in_progress", "planned"
                        "milestones": [
                            "요구사항 분석 완료",
                            "인프라 구축",
                            "파일럿 개발",
                        ]
                    },
                    ...
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        phases = phases or []
        if not phases:
            return slide

        num = len(phases)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        phase_w = (total_w - int(gap) * (num - 1)) // num
        chevron_h = Cm(1.5)
        status_colors = {
            "completed": LGColors.CHARCOAL,
            "in_progress": LGColors.RED,
            "planned": LGColors.MEDIUM_GRAY,
        }

        # Phase chevrons
        x = Cm(1.5)
        for pi, phase in enumerate(phases):
            status = phase.get("status", "planned")
            bg = status_colors.get(status, LGColors.MEDIUM_GRAY)

            # Chevron header
            chevron = self.add_box(
                slide, int(x), int(y_pos), phase_w, int(chevron_h),
                text=phase.get("name", ""),
                bg_color=bg, text_color=LGColors.WHITE,
                font_size=Pt(10), bold=True, shadow=False
            )

            # Period label below chevron
            period_y = int(y_pos) + int(chevron_h) + int(Cm(0.1))
            self._add_textbox(
                slide, int(x), period_y, phase_w, Cm(0.6),
                text=phase.get("period", ""),
                size=Pt(9), color=bg,
                alignment=PP_ALIGN.CENTER
            )

            # Status badge
            badge_y = period_y + int(Cm(0.7))
            status_text = {"completed": "완료", "in_progress": "진행중",
                          "planned": "예정"}.get(status, "")
            status_badge_bg = {
                "completed": LGColors.GREEN,
                "in_progress": LGColors.RED,
                "planned": LGColors.BORDER_GRAY,
            }.get(status, LGColors.BORDER_GRAY)
            self.add_box(
                slide, int(x) + phase_w // 4, badge_y,
                phase_w // 2, Cm(0.5),
                text=status_text, bg_color=status_badge_bg,
                text_color=LGColors.WHITE,
                font_size=Pt(8), bold=True, shadow=False
            )

            # Milestones below
            milestones_y = badge_y + int(Cm(0.8))
            milestones = phase.get("milestones", [])
            if milestones:
                ms_box = slide.shapes.add_textbox(
                    int(x), milestones_y, phase_w,
                    int(LGDimensions.SLIDE_HEIGHT) - milestones_y - int(Cm(1.0))
                )
                tf = ms_box.text_frame
                tf.word_wrap = True
                ms_box.fill.solid()
                ms_box.fill.fore_color.rgb = LGColors.LIGHT_GRAY
                ms_box.line.fill.background()
                tf.margin_left = Cm(0.3)
                tf.margin_right = Cm(0.2)
                tf.margin_top = Cm(0.3)

                for j, ms in enumerate(milestones):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '✓' if status == "completed" else '▸')
                    pPr.append(buChar)
                    pPr.set('marL', str(int(Cm(0.5))))
                    pPr.set('indent', str(int(Cm(-0.4))))
                    run = p.add_run()
                    run.text = ms
                    color = LGColors.BLACK if status != "planned" else LGColors.MEDIUM_GRAY
                    self._set_font(run, size=Pt(10), color=color)

                self._add_shadow(ms_box)

            x = int(x) + phase_w + int(gap)

        return slide

    def add_comparison_cards(self, title, section="", subtitle="",
                             cards=None):
        """
        Add a comparison cards slide — side-by-side cards for comparing
        options, vendors, solutions, or before/after states. Each card
        has a header, key metrics, and detail items.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            cards: List of card dicts (2-4 cards):
                [
                    {
                        "header": "Option A",
                        "header_color": "#3C3C3C",  # optional
                        "highlight": True,  # optional: mark as recommended
                        "metrics": [
                            {"label": "비용", "value": "5억원"},
                            {"label": "기간", "value": "6개월"},
                        ],
                        "pros": ["높은 확장성", "기존 시스템 호환"],
                        "cons": ["초기 비용 높음"],
                    },
                    ...
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        cards = cards or []
        if not cards:
            return slide

        num = len(cards)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        card_w = (total_w - int(gap) * (num - 1)) // num
        card_h = int(LGDimensions.SLIDE_HEIGHT) - int(y_pos) - int(Cm(0.5))
        header_h = Cm(1.0)

        x = Cm(1.5)
        for card in cards:
            header = card.get("header", "")
            highlight = card.get("highlight", False)
            metrics = card.get("metrics", [])
            pros = card.get("pros", [])
            cons = card.get("cons", [])

            hdr_color_hex = card.get("header_color", None)
            if hdr_color_hex:
                hdr_bg = RGBColor(
                    int(hdr_color_hex[1:3], 16),
                    int(hdr_color_hex[3:5], 16),
                    int(hdr_color_hex[5:7], 16)
                )
            else:
                hdr_bg = LGColors.RED if highlight else LGColors.CHARCOAL

            # Header with optional "추천" badge
            hdr_text = f"★ {header} (추천)" if highlight else header
            self.add_box(
                slide, int(x), int(y_pos), card_w, int(header_h),
                text=hdr_text, bg_color=hdr_bg,
                text_color=LGColors.WHITE,
                font_size=Pt(11), bold=True, shadow=False
            )

            # Card body
            body_top = int(y_pos) + int(header_h)
            body_h = card_h - int(header_h)
            border = LGColors.RED if highlight else LGColors.BORDER_GRAY

            body = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(x), body_top, card_w, body_h
            )
            body.fill.solid()
            body.fill.fore_color.rgb = LGColors.WHITE
            body.line.color.rgb = border
            body.line.width = Pt(1.5 if highlight else 0.75)

            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.4)
            tf.margin_right = Cm(0.4)
            tf.margin_top = Cm(0.3)

            # Metrics
            for mi, metric in enumerate(metrics):
                p = tf.paragraphs[0] if mi == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = f"{metric['label']}: "
                self._set_font(run, size=Pt(10), color=LGColors.MEDIUM_GRAY)
                run2 = p.add_run()
                run2.text = metric['value']
                self._set_font(run2, size=Pt(10), bold=True, color=LGColors.BLACK)

            # Separator
            if metrics and (pros or cons):
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                run = p.add_run()
                run.text = "─" * 20
                self._set_font(run, size=Pt(6), color=LGColors.BORDER_GRAY)

            # Pros
            if pros:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = "장점"
                self._set_font(run, size=Pt(9), bold=True, color=LGColors.GREEN)
                for item in pros:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(2)
                    run = p.add_run()
                    run.text = f"  + {item}"
                    self._set_font(run, size=Pt(9), color=LGColors.BLACK)

            # Cons
            if cons:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = "단점"
                self._set_font(run, size=Pt(9), bold=True,
                              color=RGBColor(0xA5, 0x00, 0x34))
                for item in cons:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(2)
                    run = p.add_run()
                    run.text = f"  - {item}"
                    self._set_font(run, size=Pt(9), color=LGColors.BLACK)

            self._add_shadow(body)
            x = int(x) + card_w + int(gap)

        return slide

    # ─────────────────────────────────────────
    # New Slide Types
    # ─────────────────────────────────────────

    def add_gantt_chart(self, title, section="", subtitle="", tasks=None,
                        start_date="", months=6):
        """
        Add a Gantt chart slide showing project task timelines.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            tasks: List of task dicts:
                [
                    {"name": "요구사항 분석", "start": 0, "duration": 2,
                     "progress": 100, "color": "#A50034"},
                    {"name": "설계", "start": 1, "duration": 3,
                     "progress": 60},
                    ...
                ]
                - start: month offset from start (0-based)
                - duration: number of months
                - progress: completion percentage (0-100)
                - color: optional bar color hex
            start_date: Label for the start period (e.g., "2025.01")
            months: Total number of months to display (default 6)
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        tasks = tasks or []
        if not tasks:
            return slide

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Layout dimensions
        name_col_width = Cm(6)
        chart_left = int(Cm(1.5)) + int(name_col_width)
        chart_width = sw - chart_left - int(Cm(0.5))
        month_width = chart_width // months
        row_height = Cm(1.0)
        bar_height = Cm(0.6)
        header_height = Cm(0.8)

        # Month header row
        for m in range(months):
            mx = chart_left + month_width * m
            # Header cell
            self.add_box(
                slide, mx, int(y_pos), month_width, int(header_height),
                text=f"M{m + 1}" if not start_date else "",
                bg_color=LGColors.CHARCOAL, text_color=LGColors.WHITE,
                font_size=Pt(9), bold=True, shadow=False
            )
            # Vertical grid line
            if m > 0:
                grid_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    mx, int(y_pos) + int(header_height),
                    Pt(0.5), int(row_height) * len(tasks)
                )
                grid_line.fill.solid()
                grid_line.fill.fore_color.rgb = LGColors.BORDER_GRAY
                grid_line.line.fill.background()

        # Month labels if start_date provided
        if start_date:
            try:
                base_year = int(start_date.split('.')[0])
                base_month = int(start_date.split('.')[1])
            except (ValueError, IndexError):
                base_year, base_month = 2025, 1
            for m in range(months):
                cm = ((base_month - 1 + m) % 12) + 1
                cy = base_year + (base_month - 1 + m) // 12
                mx = chart_left + month_width * m
                self._add_textbox(
                    slide, mx, int(y_pos),
                    month_width, int(header_height),
                    text=f"{cy}.{cm:02d}",
                    size=Pt(8), bold=True, color=LGColors.WHITE,
                    alignment=PP_ALIGN.CENTER
                )

        # Task name header
        self.add_box(
            slide, Cm(1.5), int(y_pos), int(name_col_width), int(header_height),
            text="작업 항목", bg_color=LGColors.CHARCOAL,
            text_color=LGColors.WHITE, font_size=Pt(9), bold=True, shadow=False
        )

        y_pos = int(y_pos) + int(header_height)

        # Task rows
        for ti, task in enumerate(tasks):
            task_name = task.get("name", "")
            t_start = task.get("start", 0)
            t_dur = task.get("duration", 1)
            t_progress = task.get("progress", 0)
            t_color_hex = task.get("color", "#A50034")

            try:
                tr, tg, tb = int(t_color_hex[1:3], 16), int(t_color_hex[3:5], 16), int(t_color_hex[5:7], 16)
                bar_color = RGBColor(tr, tg, tb)
            except (ValueError, IndexError):
                bar_color = LGColors.RED

            row_y = y_pos + int(row_height) * ti

            # Alternating row background
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                chart_left, row_y, chart_width, int(row_height)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = LGColors.WHITE if ti % 2 == 0 else LGColors.LIGHT_GRAY
            row_bg.line.fill.background()

            # Task name
            self._add_textbox(
                slide, Cm(1.5), row_y,
                int(name_col_width), int(row_height),
                text=task_name, size=Pt(10), color=LGColors.BLACK,
                alignment=PP_ALIGN.LEFT,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Bar background (total duration)
            bar_x = chart_left + month_width * t_start
            bar_w = month_width * t_dur
            bar_y = row_y + (int(row_height) - int(bar_height)) // 2

            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_y, bar_w, int(bar_height)
            )
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            bg_bar.line.fill.background()

            # Progress fill
            if t_progress > 0:
                fill_w = max(int(bar_w * t_progress / 100), int(Cm(0.1)))
                fill_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bar_x, bar_y, fill_w, int(bar_height)
                )
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = bar_color
                fill_bar.line.fill.background()

            # Progress label
            if t_progress > 0:
                self._add_textbox(
                    slide, bar_x, bar_y,
                    bar_w, int(bar_height),
                    text=f"{t_progress}%",
                    size=Pt(7), bold=True, color=LGColors.WHITE,
                    alignment=PP_ALIGN.CENTER,
                    vertical=MSO_ANCHOR.MIDDLE
                )

        return slide

    def add_org_chart(self, title, section="", subtitle="", org_data=None):
        """
        Add an organizational chart slide with hierarchical boxes and connectors.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            org_data: Dict representing the org hierarchy:
                {
                    "name": "CEO",
                    "title": "대표이사",
                    "children": [
                        {
                            "name": "CTO",
                            "title": "기술본부장",
                            "children": [
                                {"name": "Dev Lead", "title": "개발팀장"},
                                {"name": "QA Lead", "title": "품질팀장"},
                            ]
                        },
                        {
                            "name": "CFO",
                            "title": "재무본부장",
                        }
                    ]
                }
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        if not org_data:
            return slide

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Calculate tree structure
        def count_leaves(node):
            children = node.get("children", [])
            if not children:
                return 1
            return sum(count_leaves(c) for c in children)

        def get_depth(node):
            children = node.get("children", [])
            if not children:
                return 1
            return 1 + max(get_depth(c) for c in children)

        depth = get_depth(org_data)
        total_leaves = count_leaves(org_data)

        # Box dimensions
        box_w = Cm(4.5)
        box_h = Cm(1.6)
        h_gap = Cm(0.5)
        v_gap = Cm(1.2)
        level_height = int(box_h) + int(v_gap)

        available_width = sw - int(Cm(2.0))
        content_top = int(y_pos)

        # Recursive drawing with position tracking
        def draw_node(node, level, cx, is_root=False):
            """Draw a node at center_x=cx, returns (center_x, top_y, bottom_y)."""
            ny = content_top + level * level_height
            nx = cx - int(box_w) // 2

            # Box styling
            bg = LGColors.CHARCOAL if is_root else LGColors.LIGHT_GRAY
            tc = LGColors.WHITE if is_root else LGColors.BLACK

            box = self.add_box(
                slide, nx, ny, int(box_w), int(box_h),
                bg_color=bg, shadow=True
            )

            # Name (bold)
            name = node.get("name", "")
            node_title = node.get("title", "")
            name_box = slide.shapes.add_textbox(
                nx, ny, int(box_w), int(box_h) // 2
            )
            name_tf = name_box.text_frame
            name_tf.word_wrap = True
            name_box.fill.background()
            name_box.line.fill.background()
            p = name_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = name
            self._set_font(run, size=Pt(10), bold=True, color=tc)

            # Title
            if node_title:
                title_box = slide.shapes.add_textbox(
                    nx, ny + int(box_h) // 2, int(box_w), int(box_h) // 2
                )
                title_tf = title_box.text_frame
                title_tf.word_wrap = True
                title_box.fill.background()
                title_box.line.fill.background()
                p = title_tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = node_title
                self._set_font(run, size=Pt(8), color=tc if is_root else LGColors.MEDIUM_GRAY)

            children = node.get("children", [])
            if children:
                # Calculate children positions
                child_leaves = [count_leaves(c) for c in children]
                total_child_leaves = sum(child_leaves)

                # Minimum width per leaf
                min_leaf_w = int(box_w) + int(h_gap)
                total_children_width = total_child_leaves * min_leaf_w

                # Center children under parent
                child_start_x = cx - total_children_width // 2
                child_centers = []

                offset = 0
                for ci, child in enumerate(children):
                    cl = child_leaves[ci]
                    child_cx = child_start_x + offset + (cl * min_leaf_w) // 2
                    child_centers.append(child_cx)
                    draw_node(child, level + 1, child_cx)
                    offset += cl * min_leaf_w

                # Draw connectors
                parent_bottom_y = ny + int(box_h)
                child_top_y = content_top + (level + 1) * level_height
                mid_y = (parent_bottom_y + child_top_y) // 2

                # Vertical line from parent down
                v_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    cx - Pt(1), parent_bottom_y,
                    Pt(2), mid_y - parent_bottom_y
                )
                v_line.fill.solid()
                v_line.fill.fore_color.rgb = LGColors.BORDER_GRAY
                v_line.line.fill.background()

                # Horizontal line connecting children
                if len(child_centers) > 1:
                    h_left = min(child_centers)
                    h_right = max(child_centers)
                    h_line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        h_left, mid_y,
                        h_right - h_left, Pt(2)
                    )
                    h_line.fill.solid()
                    h_line.fill.fore_color.rgb = LGColors.BORDER_GRAY
                    h_line.line.fill.background()

                # Vertical lines from horizontal to each child
                for ccx in child_centers:
                    vc_line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        ccx - Pt(1), mid_y,
                        Pt(2), child_top_y - mid_y
                    )
                    vc_line.fill.solid()
                    vc_line.fill.fore_color.rgb = LGColors.BORDER_GRAY
                    vc_line.line.fill.background()

        # Draw from root centered on slide
        draw_node(org_data, 0, sw // 2, is_root=True)

        return slide

    def add_pyramid(self, title, section="", subtitle="", levels=None):
        """
        Add a pyramid/hierarchy diagram slide with trapezoid layers.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            levels: List of level dicts (top to bottom):
                [
                    {"label": "전략", "description": "장기 비전 및 목표"},
                    {"label": "전술", "description": "실행 계획"},
                    {"label": "운영", "description": "일상 업무 프로세스"},
                ]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        levels = levels or []
        if not levels:
            return slide

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        num_levels = len(levels)
        available_height = sh - int(y_pos) - int(Cm(0.5))
        level_height = available_height // num_levels
        gap = Cm(0.1)

        # Pyramid center and width range
        center_x = sw // 2 - int(Cm(2))  # Shift left to leave room for descriptions
        min_width = Cm(6)
        max_width = Cm(18)
        desc_left = center_x + int(max_width) // 2 + int(Cm(1.0))

        # Distinct color palette for each level
        _pyramid_palette = [
            "#A50034", "#3C3C3C", "#1565C0", "#2E7D32", "#D4760A", "#7B1FA2"
        ]
        for i, level in enumerate(levels):
            # Calculate width for this level (narrower at top)
            fraction = i / max(num_levels - 1, 1)
            level_w = int(min_width) + int((int(max_width) - int(min_width)) * fraction)
            lx = center_x - level_w // 2
            ly = int(y_pos) + level_height * i

            # Color: use per-level color key, or palette, or fallback
            level_color_hex = level.get("color", None)
            if not level_color_hex:
                level_color_hex = _pyramid_palette[i % len(_pyramid_palette)]
            level_color = RGBColor(
                int(level_color_hex[1:3], 16),
                int(level_color_hex[3:5], 16),
                int(level_color_hex[5:7], 16)
            )

            # Trapezoid shape (approximated with rectangle + rounded corners)
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                lx, ly, level_w, level_height - int(gap)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = level_color
            box.line.fill.background()

            # Level label (centered in trapezoid)
            label = level.get("label", "")
            self._add_textbox(
                slide, lx, ly,
                level_w, level_height - int(gap),
                text=label, size=Pt(14), bold=True,
                color=LGColors.WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Description on the right
            desc = level.get("description", "")
            if desc:
                # Connector line
                conn_y = ly + (level_height - int(gap)) // 2
                conn = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    lx + level_w, conn_y,
                    desc_left - (lx + level_w), Pt(1)
                )
                conn.fill.solid()
                conn.fill.fore_color.rgb = LGColors.BORDER_GRAY
                conn.line.fill.background()

                # Description text
                self._add_textbox(
                    slide, desc_left, ly,
                    sw - desc_left - int(Cm(0.5)), level_height - int(gap),
                    text=desc, size=Pt(10), color=LGColors.DARK_GRAY,
                    alignment=PP_ALIGN.LEFT,
                    vertical=MSO_ANCHOR.MIDDLE
                )

        return slide

    def add_positioning_map(self, title, section="", subtitle="",
                            x_label="", y_label="", items=None,
                            quadrant_labels=None):
        """
        Add a 2D positioning/perceptual map slide with plotted items.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            x_label: Label for X axis
            y_label: Label for Y axis
            items: List of item dicts:
                [
                    {"name": "제품A", "x": 0.8, "y": 0.7,
                     "size": "large", "color": "#A50034"},
                    {"name": "제품B", "x": 0.3, "y": 0.4,
                     "size": "medium"},
                    ...
                ]
                - x, y: normalized position (0.0 to 1.0)
                - size: "small", "medium", "large" (default "medium")
                - color: optional hex color
            quadrant_labels: Optional list of 4 labels:
                ["고가치/고성장", "고가치/저성장",
                 "저가치/고성장", "저가치/저성장"]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        items = items or []

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Map area
        map_left = int(Cm(3.5))
        map_top = int(y_pos) + int(Cm(0.5))
        map_size = min(sw - int(Cm(5.0)), sh - map_top - int(Cm(1.5)))
        map_width = map_size
        map_height = map_size

        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            map_left, map_top, map_width, map_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = LGColors.LIGHT_GRAY
        bg.line.color.rgb = LGColors.BORDER_GRAY
        bg.line.width = Pt(1)

        # Cross-hairs (axes through center)
        mid_x = map_left + map_width // 2
        mid_y = map_top + map_height // 2

        # Vertical axis
        v_axis = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            mid_x - Pt(1), map_top, Pt(2), map_height
        )
        v_axis.fill.solid()
        v_axis.fill.fore_color.rgb = LGColors.BORDER_GRAY
        v_axis.line.fill.background()

        # Horizontal axis
        h_axis = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            map_left, mid_y - Pt(1), map_width, Pt(2)
        )
        h_axis.fill.solid()
        h_axis.fill.fore_color.rgb = LGColors.BORDER_GRAY
        h_axis.line.fill.background()

        # Axis labels
        if x_label:
            self._add_textbox(
                slide, map_left, map_top + map_height + int(Cm(0.2)),
                map_width, Cm(0.8),
                text=x_label, size=Pt(10), bold=True,
                color=LGColors.DARK_GRAY, alignment=PP_ALIGN.CENTER
            )
        if y_label:
            self._add_textbox(
                slide, map_left - int(Cm(3.0)), map_top,
                Cm(2.5), map_height,
                text=y_label, size=Pt(10), bold=True,
                color=LGColors.DARK_GRAY, alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # Quadrant labels
        if quadrant_labels and len(quadrant_labels) >= 4:
            ql_positions = [
                (map_left + int(Cm(0.3)), map_top + int(Cm(0.2))),
                (mid_x + int(Cm(0.3)), map_top + int(Cm(0.2))),
                (map_left + int(Cm(0.3)), mid_y + int(Cm(0.2))),
                (mid_x + int(Cm(0.3)), mid_y + int(Cm(0.2))),
            ]
            for qi, (qx, qy) in enumerate(ql_positions):
                if qi < len(quadrant_labels):
                    self._add_textbox(
                        slide, qx, qy,
                        map_width // 2 - int(Cm(0.6)), Cm(0.6),
                        text=quadrant_labels[qi],
                        size=Pt(8), color=LGColors.MEDIUM_GRAY,
                        alignment=PP_ALIGN.LEFT
                    )

        # Plot items
        size_map = {"small": Cm(1.2), "medium": Cm(1.8), "large": Cm(2.5)}
        for item in items:
            ix = item.get("x", 0.5)
            iy = item.get("y", 0.5)
            i_size = size_map.get(item.get("size", "medium"), Cm(1.8))
            i_name = item.get("name", "")
            i_color_hex = item.get("color", "#A50034")

            try:
                ir, ig, ib = int(i_color_hex[1:3], 16), int(i_color_hex[3:5], 16), int(i_color_hex[5:7], 16)
                i_color = RGBColor(ir, ig, ib)
            except (ValueError, IndexError):
                i_color = LGColors.RED

            # Convert normalized coords to map coords
            px = map_left + int(map_width * ix) - int(i_size) // 2
            py = map_top + int(map_height * (1.0 - iy)) - int(i_size) // 2

            # Bubble
            bubble = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                px, py, int(i_size), int(i_size)
            )
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = i_color

            # Set transparency (40%)
            spPr = bubble._element.spPr
            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    alpha = _make_oxml_element('a:alpha')
                    alpha.set('val', '60000')  # 60% opaque
                    srgbClr.append(alpha)

            bubble.line.fill.background()

            # Label
            self._add_textbox(
                slide, px, py,
                int(i_size), int(i_size),
                text=i_name, size=Pt(8), bold=True,
                color=LGColors.WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        return slide

    def add_keyword_highlight(self, title, section="", subtitle="",
                              keywords=None, description=""):
        """
        Add a keyword emphasis/tag-cloud style slide.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            keywords: List of keyword dicts:
                [
                    {"text": "AI 혁신", "size": 36, "color": "#A50034"},
                    {"text": "디지털 전환", "size": 28},
                    {"text": "데이터 분석", "size": 24, "color": "#2E7D32"},
                    ...
                ]
                - size: font size in pt (determines visual weight)
                - color: optional hex color (default LG RED)
            description: Supporting description text below keywords
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        keywords = keywords or []
        if not keywords:
            return slide

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        # Keyword cloud area
        cloud_left = Cm(2.0)
        cloud_top = int(y_pos) + int(Cm(0.5))
        cloud_width = sw - int(Cm(4.0))
        cloud_height = sh - cloud_top - int(Cm(3.0)) if description else sh - cloud_top - int(Cm(1.0))

        # Arrange keywords in rows
        current_x = int(cloud_left)
        current_y = cloud_top
        row_height = 0
        padding = Cm(0.3)

        for kw in keywords:
            kw_text = kw.get("text", "")
            kw_size = kw.get("size", 24)
            kw_color_hex = kw.get("color", "#A50034")

            try:
                kr, kg, kb = int(kw_color_hex[1:3], 16), int(kw_color_hex[3:5], 16), int(kw_color_hex[5:7], 16)
                kw_color = RGBColor(kr, kg, kb)
            except (ValueError, IndexError):
                kw_color = LGColors.RED

            # Estimate box size based on text length and font size
            char_width = int(Pt(kw_size)) * 0.7
            box_w = int(len(kw_text) * char_width + int(Cm(1.0)))
            box_h = int(Pt(kw_size)) * 2 + int(Cm(0.3))

            # Wrap to next row if needed
            if current_x + box_w > int(cloud_left) + cloud_width:
                current_x = int(cloud_left)
                current_y += row_height + int(padding)
                row_height = 0

            row_height = max(row_height, box_h)

            # Background box with slight opacity
            bg_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                current_x, current_y, box_w, box_h
            )
            bg_box.fill.solid()
            bg_box.fill.fore_color.rgb = LGColors.LIGHT_GRAY
            bg_box.line.fill.background()

            # Keyword text
            self._add_textbox(
                slide, current_x, current_y,
                box_w, box_h,
                text=kw_text, size=Pt(kw_size), bold=True,
                color=kw_color,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            current_x += box_w + int(padding)

        # Description below
        if description:
            desc_top = sh - int(Cm(2.5))
            self._add_textbox(
                slide, Cm(2.0), desc_top,
                sw - int(Cm(4.0)), Cm(2.0),
                text=description, size=LGTypography.BODY,
                color=LGColors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    # ─────────────────────────────────────────
    # Swimlane Diagram
    # ─────────────────────────────────────────

    def add_swimlane(self, title, section="", subtitle="",
                     lanes=None, steps=None, connections=None):
        """
        Add a swimlane process diagram — horizontal lanes for each
        role/department/system with process steps placed in columns.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            lanes: List of lane names (roles/departments):
                ["고객사", "PM", "개발팀", "QA팀"]
            steps: List of step dicts:
                [
                    {"lane": 0, "col": 0, "text": "요구사항 전달"},
                    {"lane": 1, "col": 1, "text": "요구사항 분석"},
                    {"lane": 2, "col": 2, "text": "설계/개발", "color": "#1565C0"},
                    {"lane": 3, "col": 3, "text": "테스트"},
                    {"lane": 1, "col": 4, "text": "검수 요청"},
                    {"lane": 0, "col": 5, "text": "최종 승인"},
                ]
                - lane: index into lanes list (which row)
                - col: column position (left to right sequence)
                - text: step label
                - color: optional background color hex
            connections: List of (from_step_index, to_step_index) tuples:
                [(0,1), (1,2), (2,3), (3,4), (4,5)]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        lanes = lanes or []
        steps = steps or []
        connections = connections or []
        if not lanes or not steps:
            return slide

        y_start = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_start, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_start = Cm(2.8)
        else:
            y_start = Cm(2.2)

        num_lanes = len(lanes)
        num_cols = max((s.get("col", 0) for s in steps), default=0) + 1

        # Layout dimensions
        label_width = Cm(3.0)
        grid_left = int(Cm(1.5)) + int(label_width) + int(Cm(0.2))
        grid_right_margin = Cm(1.0)
        grid_width = sw - grid_left - int(grid_right_margin)
        available_height = sh - int(y_start) - int(Cm(0.5))
        lane_height = available_height // num_lanes
        col_width = grid_width // num_cols

        # Step box dimensions
        step_w = int(col_width * 0.75)
        step_h = int(lane_height * 0.55)

        # Distinct lane colors (alternating light backgrounds)
        lane_bg_colors = [
            RGBColor(0xF2, 0xF2, 0xF2),  # Light Gray
            RGBColor(0xFF, 0xFF, 0xFF),   # White
        ]

        # Lane label accent colors
        lane_accent_colors = [
            LGColors.RED,
            LGColors.CHARCOAL,
            RGBColor(0x15, 0x65, 0xC0),  # Blue
            RGBColor(0x2E, 0x7D, 0x32),  # Green
            RGBColor(0xD4, 0x76, 0x0A),  # Orange
            RGBColor(0x7B, 0x1F, 0xA2),  # Purple
        ]

        # Draw lanes (horizontal bands)
        for i, lane_name in enumerate(lanes):
            ly = int(y_start) + lane_height * i

            # Lane background band
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                grid_left, ly, grid_width, lane_height
            )
            band.fill.solid()
            band.fill.fore_color.rgb = lane_bg_colors[i % 2]
            band.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            band.line.width = Pt(0.5)

            # Lane label (left side)
            accent_color = lane_accent_colors[i % len(lane_accent_colors)]
            label_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(Cm(1.5)), ly, int(label_width), lane_height
            )
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = accent_color
            label_box.line.fill.background()

            tf = label_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.2)
            tf.margin_right = Cm(0.2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = lane_name
            self._set_font(run, size=Pt(11), bold=True, color=LGColors.WHITE)

            # Thin horizontal divider at bottom of lane
            if i < num_lanes - 1:
                div_y = ly + lane_height
                div = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    grid_left, div_y, grid_width, Pt(1)
                )
                div.fill.solid()
                div.fill.fore_color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
                div.line.fill.background()

        # Place step boxes and record centers for connections
        step_centers = {}  # step_index -> (cx, cy)

        for idx, step in enumerate(steps):
            lane_idx = step.get("lane", 0)
            col_idx = step.get("col", 0)
            text = step.get("text", "")
            color_hex = step.get("color", None)

            # Calculate position (centered in cell)
            cx = grid_left + col_width * col_idx + col_width // 2
            cy = int(y_start) + lane_height * lane_idx + lane_height // 2
            bx = cx - step_w // 2
            by = cy - step_h // 2

            # Step box color
            if color_hex:
                box_color = RGBColor(
                    int(color_hex[1:3], 16),
                    int(color_hex[3:5], 16),
                    int(color_hex[5:7], 16)
                )
            else:
                box_color = LGColors.CHARCOAL

            # Step box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bx, by, step_w, step_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = box_color
            box.line.fill.background()

            # Adjust corner radius
            try:
                box.adjustments[0] = 0.15
            except Exception:
                pass

            # Step text
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.15)
            tf.margin_right = Cm(0.15)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run, size=Pt(10), bold=True, color=LGColors.WHITE)

            step_centers[idx] = (cx, cy)

            self._add_shadow(box)

        # Draw connection arrows
        for from_idx, to_idx in connections:
            if from_idx not in step_centers or to_idx not in step_centers:
                continue

            fx, fy = step_centers[from_idx]
            tx, ty = step_centers[to_idx]

            from_step = steps[from_idx]
            to_step = steps[to_idx]
            from_col = from_step.get("col", 0)
            to_col = to_step.get("col", 0)
            from_lane = from_step.get("lane", 0)
            to_lane = to_step.get("lane", 0)

            if from_col == to_col:
                # Vertical arrow (same column, different lane)
                if fy < ty:
                    ay = fy + step_h // 2
                    ah = ty - step_h // 2 - ay
                else:
                    ay = ty + step_h // 2
                    ah = fy - step_h // 2 - ay

                if ah > 0:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW if fy < ty else MSO_SHAPE.UP_ARROW,
                        fx - int(Cm(0.2)), ay,
                        int(Cm(0.4)), ah
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = LGColors.CHARCOAL
                    arrow.line.fill.background()
            else:
                # Horizontal or diagonal → use right arrow
                ax = fx + step_w // 2
                aw = tx - step_w // 2 - ax

                if from_lane == to_lane and aw > 0:
                    # Same lane: horizontal arrow
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        ax, fy - int(Cm(0.25)),
                        aw, int(Cm(0.5))
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    arrow.line.fill.background()
                elif aw > 0:
                    # Different lane: L-shaped connector (horizontal then vertical)
                    mid_x = ax + aw // 2

                    # Horizontal segment
                    h_seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        ax, fy - Pt(2),
                        mid_x - ax + int(Cm(0.1)), Pt(4)
                    )
                    h_seg.fill.solid()
                    h_seg.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    h_seg.line.fill.background()

                    # Vertical segment
                    v_top = min(fy, ty)
                    v_bottom = max(fy, ty)
                    v_seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        mid_x, v_top,
                        Pt(4), v_bottom - v_top
                    )
                    v_seg.fill.solid()
                    v_seg.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    v_seg.line.fill.background()

                    # Final horizontal to target
                    h_seg2 = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        mid_x, ty - Pt(2),
                        tx - step_w // 2 - mid_x, Pt(4)
                    )
                    h_seg2.fill.solid()
                    h_seg2.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    h_seg2.line.fill.background()

                    # Arrow head at target
                    arr_head = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        tx - step_w // 2 - int(Cm(0.5)), ty - int(Cm(0.25)),
                        int(Cm(0.5)), int(Cm(0.5))
                    )
                    arr_head.fill.solid()
                    arr_head.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    arr_head.line.fill.background()

        return slide

    # ─────────────────────────────────────────
    # Chart / Image Slides
    # ─────────────────────────────────────────

    def add_chart_slide(self, title, section="", subtitle="",
                        chart_path=None, caption=""):
        """
        Add a slide with a chart image (from matplotlib save or any image file).

        Usage:
            import matplotlib.pyplot as plt
            fig, ax = plt.subplots()
            ax.bar(['A', 'B', 'C'], [10, 20, 15])
            fig.savefig('/tmp/chart.png', dpi=150, bbox_inches='tight')
            plt.close()
            prs.add_chart_slide("매출 분석", chart_path='/tmp/chart.png')

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            chart_path: Path to chart image file (PNG, JPG, etc.)
            caption: Optional caption text below the chart
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        if chart_path and os.path.exists(chart_path):
            # Calculate image dimensions to fill content area
            img_top = int(y_pos) + int(Cm(0.3))
            img_left = Cm(2.0)
            img_max_width = sw - int(Cm(4.0))
            img_max_height = sh - img_top - int(Cm(1.5))
            if caption:
                img_max_height -= int(Cm(1.0))

            # Add image centered
            pic = slide.shapes.add_picture(
                chart_path,
                img_left, img_top,
                width=img_max_width
            )

            # Scale to fit height if needed
            if pic.height > img_max_height:
                ratio = img_max_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = img_max_height

            # Center horizontally
            pic.left = (sw - pic.width) // 2

        if caption:
            cap_top = sh - int(Cm(1.5))
            self._add_textbox(
                slide, Cm(2.0), cap_top,
                sw - int(Cm(4.0)), Cm(1.0),
                text=caption, size=LGTypography.CAPTION,
                color=LGColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    def add_image_slide(self, title, section="", subtitle="", images=None):
        """
        Add a slide with one or more images in various layouts.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            images: List of image dicts:
                [
                    {"path": "/path/to/img.png", "width": 15, "caption": "설명"},
                    ...
                ]
                - 1 image: centered, large
                - 2 images: side by side
                - 3+ images: grid layout
                - width: optional width in cm
                - caption: optional caption text
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        self._add_accent_bar(slide)
        if section:
            self._add_section_indicator(slide, section)
        self._add_slide_title(slide, title)

        images = images or []
        if not images:
            return slide

        y_pos = Cm(1.8)
        if subtitle:
            self._add_textbox(
                slide, Cm(1.5), y_pos, Cm(28), Cm(1.0),
                text=subtitle, size=LGTypography.SUBTITLE,
                color=LGColors.DARK_GRAY, weight="semibold"
            )
            y_pos = Cm(2.8)
        else:
            y_pos = Cm(2.2)

        content_top = int(y_pos) + int(Cm(0.3))
        content_height = sh - content_top - int(Cm(1.0))
        content_width = sw - int(Cm(3.0))
        margin_left = Cm(1.5)
        num_images = len(images)

        if num_images == 1:
            # Single image centered
            img = images[0]
            path = img.get("path", "")
            caption = img.get("caption", "")
            img_w = Cm(img.get("width", 20)) if img.get("width") else content_width

            if os.path.exists(path):
                pic = slide.shapes.add_picture(
                    path, int(margin_left), content_top, width=int(img_w)
                )
                # Scale to fit
                if pic.height > content_height - int(Cm(1.0)):
                    ratio = (content_height - int(Cm(1.0))) / pic.height
                    pic.width = int(pic.width * ratio)
                    pic.height = int(content_height - int(Cm(1.0)))
                # Center
                pic.left = (sw - pic.width) // 2

                if caption:
                    self._add_textbox(
                        slide, int(margin_left),
                        content_top + pic.height + int(Cm(0.2)),
                        content_width, Cm(0.8),
                        text=caption, size=LGTypography.CAPTION,
                        color=LGColors.MEDIUM_GRAY,
                        alignment=PP_ALIGN.CENTER
                    )

        elif num_images == 2:
            # Side by side
            gap = Cm(0.5)
            each_w = (content_width - int(gap)) // 2

            for i, img in enumerate(images):
                path = img.get("path", "")
                caption = img.get("caption", "")
                ix = int(margin_left) + i * (each_w + int(gap))

                if os.path.exists(path):
                    pic = slide.shapes.add_picture(
                        path, ix, content_top, width=each_w
                    )
                    if pic.height > content_height - int(Cm(1.5)):
                        ratio = (content_height - int(Cm(1.5))) / pic.height
                        pic.width = int(pic.width * ratio)
                        pic.height = int(content_height - int(Cm(1.5)))

                    if caption:
                        self._add_textbox(
                            slide, ix,
                            content_top + pic.height + int(Cm(0.2)),
                            each_w, Cm(0.8),
                            text=caption, size=LGTypography.CAPTION,
                            color=LGColors.MEDIUM_GRAY,
                            alignment=PP_ALIGN.CENTER
                        )

        else:
            # Grid layout (2 columns)
            gap = Cm(0.3)
            cols = 2
            rows_count = (num_images + cols - 1) // cols
            each_w = (content_width - int(gap) * (cols - 1)) // cols
            each_h = (content_height - int(gap) * (rows_count - 1)) // rows_count

            for i, img in enumerate(images):
                path = img.get("path", "")
                caption = img.get("caption", "")
                col = i % cols
                row = i // cols
                ix = int(margin_left) + col * (each_w + int(gap))
                iy = content_top + row * (each_h + int(gap))

                if os.path.exists(path):
                    pic = slide.shapes.add_picture(
                        path, ix, iy, width=each_w
                    )
                    if pic.height > each_h - int(Cm(1.0)):
                        ratio = (each_h - int(Cm(1.0))) / pic.height
                        pic.width = int(pic.width * ratio)
                        pic.height = int(each_h - int(Cm(1.0)))

                    if caption:
                        self._add_textbox(
                            slide, ix,
                            iy + pic.height + int(Cm(0.1)),
                            each_w, Cm(0.6),
                            text=caption, size=Pt(8),
                            color=LGColors.MEDIUM_GRAY,
                            alignment=PP_ALIGN.CENTER
                        )

        return slide

    def add_matplotlib_chart(self, title, section="", subtitle="",
                             fig=None, caption=""):
        """
        Add a matplotlib figure directly to a slide (saves to temp file internally).

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            fig: matplotlib Figure object
            caption: Optional caption text below the chart

        Usage:
            import matplotlib.pyplot as plt
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.bar(['A', 'B', 'C'], [10, 20, 15])
            prs.add_matplotlib_chart("분석 결과", fig=fig)
            plt.close(fig)
        """
        if fig is None:
            return self.add_chart_slide(title, section, subtitle, None, caption)

        import tempfile
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        tmp_path = tmp.name
        tmp.close()

        try:
            fig.savefig(tmp_path, dpi=150, bbox_inches='tight',
                        facecolor='white', edgecolor='none')
            slide = self.add_chart_slide(title, section, subtitle,
                                         tmp_path, caption)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

        return slide

    def save(self, filename):
        """Save the presentation to a file."""
        self.prs.save(filename)
        return filename
